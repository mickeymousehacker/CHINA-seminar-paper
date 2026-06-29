"""
FINAL unified presentation — "China Strategies in the Automotive Industry"
Seminar: M&A as Transformer of the Automotive Industry · Goethe University Frankfurt
Authors: Julia Ort · Batughan Maus · Chiara Stubner · 29 June 2026

One consistent design system (Cambria headlines / Calibri body, China-EV navy palette,
13.33x7.5 widescreen). 21 slides, paper-matched, no content repetition.
Speaker notes in simple B2 English, paced for a ~40-minute talk.

Run:  uv run --with python-pptx python research/presentation/build_final_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x17, 0x2A)
NAVY2  = RGBColor(0x1E, 0x29, 0x3B)
BLUE   = RGBColor(0x25, 0x63, 0xEB)   # BYD
TEAL   = RGBColor(0x06, 0xB6, 0xD4)   # Geely
RED    = RGBColor(0xDC, 0x26, 0x26)   # SAIC / highlight
BG     = RGBColor(0xF8, 0xFA, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x0F, 0x17, 0x2A)
GREY   = RGBColor(0x64, 0x74, 0x8B)
LINE   = RGBColor(0xCB, 0xD5, 0xE1)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xD9, 0x77, 0x06)

SAIC, BYD_C, GEELY = RED, BLUE, TEAL
HEAD, BODY = "Arial", "Arial"
W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)


# ── Core helpers ──────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(s, l, t, w, h, color, line_color=None, line_w=Pt(0.75)):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color; sp.line.width = line_w
    return sp


def oval(s, l, t, d, color):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    return sp


def _no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def text(s, l, t, w, h, lines, size, color, *, font=BODY, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.0, sa=2):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = lines.split("\n")
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.line_spacing = ls; p.space_before = Pt(0); p.space_after = Pt(sa)
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


def bullets(s, l, t, w, h, items, size, color=INK, *, gap=6, marker="–",
            mcolor=None, ls=1.05):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = ls
        _no_bullet(p)
        if marker:
            rm = p.add_run(); rm.text = marker + "  "
            rm.font.name = BODY; rm.font.size = Pt(size); rm.font.color.rgb = mcolor or BLUE
        r = p.add_run(); r.text = it
        r.font.name = BODY; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip()


def kicker_headline(s, kicker, headline, hsize=21):
    rect(s, 0, 0, W, H, BG)
    text(s, MARGIN, Inches(0.42), Inches(12.1), Inches(0.3), kicker.upper(), 14,
         TEAL, bold=True)
    text(s, MARGIN, Inches(0.72), Inches(12.1), Inches(1.0), headline, hsize, NAVY,
         font=HEAD, bold=True, ls=1.0)
    rect(s, MARGIN, Inches(1.62), Inches(0.9), Inches(0.045), BLUE)
    rect(s, Inches(1.6), Inches(1.64), Inches(11.1), Inches(0.012), LINE)


def footer(s, page):
    text(s, MARGIN, H - Inches(0.42), Inches(9.5), Inches(0.3),
         "China Strategies in the Automotive Industry  ·  Session 6", 10.5, GREY)
    text(s, W - Inches(1.5), H - Inches(0.42), Inches(0.9), Inches(0.3), str(page),
         10.5, GREY, align=PP_ALIGN.RIGHT)


def source(s, txt):
    text(s, MARGIN, H - Inches(0.72), Inches(12.1), Inches(0.3), txt, 11, GREY,
         italic=True)


def chip(s, l, t, w, label, color, h=Inches(0.46), size=12.5):
    rect(s, l, t, w, h, color)
    text(s, l, t, w, h, label, size, WHITE, bold=True, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


def card(s, l, t, w, h, top_color, title, items, *, tsize=16, bsize=13.5):
    rect(s, l, t, w, h, CARD, line_color=LINE)
    rect(s, l, t, w, Inches(0.09), top_color)
    text(s, l + Inches(0.18), t + Inches(0.22), w - Inches(0.36), Inches(0.6), title,
         tsize, top_color, bold=True, ls=0.98)
    bullets(s, l + Inches(0.18), t + Inches(0.86), w - Inches(0.36), h - Inches(1.04),
            items, bsize, gap=5, marker="–", mcolor=top_color)


def grid(s, top, left, col_ws, row_h, headers, rows, *, fs=13.5, hs=14,
         head_bg=NAVY, cell_colors=None):
    x = left
    for ci, hd in enumerate(headers):
        rect(s, x, top, col_ws[ci], row_h, head_bg)
        text(s, x + Inches(0.08), top, col_ws[ci] - Inches(0.16), row_h, hd, hs, WHITE,
             bold=True, align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, ls=0.95)
        x = Emu(x + col_ws[ci])
    for ri, row in enumerate(rows):
        y = Emu(top + row_h * (ri + 1))
        x = left
        base = WHITE if ri % 2 == 0 else RGBColor(0xEE, 0xF2, 0xF7)
        for ci, cell in enumerate(row):
            cc = base
            if cell_colors and (ri, ci) in cell_colors:
                cc = cell_colors[(ri, ci)]
            rect(s, x, y, col_ws[ci], row_h, cc, line_color=LINE, line_w=Pt(0.5))
            fg = WHITE if (cell_colors and (ri, ci) in cell_colors) else INK
            text(s, x + Inches(0.08), y, col_ws[ci] - Inches(0.16), row_h, cell, fs, fg,
                 bold=(ci == 0), align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, ls=0.95)
            x = Emu(x + col_ws[ci])


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════════
def s01_title(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, BLUE)
    text(s, MARGIN, Inches(0.6), Inches(12), Inches(0.35),
         "M&A AS TRANSFORMER OF THE AUTOMOTIVE INDUSTRY", 13, TEAL, bold=True)
    text(s, MARGIN, Inches(2.0), Inches(12), Inches(1.7),
         "China Strategies in the\nAutomotive Industry", 46, WHITE, font=HEAD,
         bold=True, ls=1.0)
    text(s, MARGIN, Inches(3.95), Inches(11.6), Inches(0.7),
         "Which governance and financing model creates long-term shareholder value?",
         18, RGBColor(0xCB, 0xD5, 0xE1), italic=True)
    x = MARGIN
    for lbl, c in [("SAIC  ·  Joint Venture", SAIC), ("BYD  ·  Vertical Integration", BYD_C),
                   ("Geely  ·  M&A-Driven", GEELY)]:
        chip(s, x, Inches(4.95), Inches(3.7), lbl, c); x = Emu(x + Inches(3.95))
    text(s, MARGIN, Inches(6.05), Inches(12), Inches(0.4),
         "Julia Ort  ·  Batughan Maus  ·  Chiara Stubner", 15, WHITE, bold=True)
    text(s, MARGIN, Inches(6.5), Inches(12), Inches(0.6),
         "Johann Wolfgang Goethe-Universität Frankfurt  ·  Chair of Banking and Finance\n"
         "Prof. Dr. Mark Wahrenburg  ·  Advisor: Edoardo Pilla  ·  29 June 2026",
         12, GREY, ls=1.1)
    notes(s, """Good afternoon, and welcome to our presentation. Our paper asks one clear
question: how do you build a successful car company in China, and which ownership and money
structure creates the most value for shareholders over time. We answer it by comparing three
real firms that each chose a different path — SAIC with joint ventures, BYD by building
everything itself, and Geely by buying companies abroad. Our main finding, which we will defend
across the whole talk, is that governance — who controls the firm — matters more than the
technology a firm starts with. We will move through four lenses: governance, financing,
capability building, and long-term survival, and then give our verdict. At the end we also
reflect openly on how we used AI in the research. Let us start with why this question is urgent
right now.""")


def s02_intro(prs):
    s = blank(prs)
    kicker_headline(s, "Introduction · Why now",
                    "The 30-year joint-venture model is breaking down as EVs take over")
    stats = [("~55%", "EV share of new-car\nsales in China (2025)", BLUE),
             ("~60%", "of global EV sales\nfrom Chinese makers", TEAL),
             ("35.3%", "EU duty on SAIC\n(17.0% BYD · 18.8% Geely)", RED),
             ("88%", "drop in SAIC net\nprofit in 2024", RED)]
    x = MARGIN
    for num, lbl, c in stats:
        rect(s, x, Inches(1.9), Inches(2.95), Inches(2.0), CARD, line_color=LINE)
        rect(s, x, Inches(1.9), Inches(2.95), Inches(0.09), c)
        text(s, x, Inches(2.25), Inches(2.95), Inches(0.9), num, 40, c, font=HEAD,
             bold=True, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.12), Inches(3.15), Inches(2.71), Inches(0.65), lbl, 12, INK,
             align=PP_ALIGN.CENTER, ls=1.0)
        x = Emu(x + Inches(3.05))
    bullets(s, MARGIN, Inches(4.3), Inches(12.1), Inches(2.0), [
        "For three decades foreign technology was traded for market access, with a state-owned partner in the middle (Holmes et al., 2015) — the EV shift makes that logic obsolete.",
        "Firms that leaned on foreign partners are losing share fast, while independent Chinese makers have become global leaders (Whitfield & Wuttke, 2026).",
        "The EU now prices governance risk directly: the tariff a firm pays tracks how state-dependent it is (European Commission, 2024).",
    ], 15, gap=9)
    source(s, "Sources: IEA Global EV Outlook 2026 · SAIC Annual Report 2024/2025 · European Commission (2024)")
    footer(s, 2)
    notes(s, """The joint-venture era is ending, and the numbers make that clear. Electric vehicles
are now more than half of all new cars sold in China, and Chinese firms make about sixty percent
of the world's EVs. For thirty years the rule was simple: foreign companies gave their technology
and a Chinese state firm gave market access. That trade worked while engines and gearboxes were
the hard part. In the electric era, the technology that matters — batteries, power electronics,
software — sits with the Chinese independents, so the old trade no longer helps. You can see the
stress at SAIC, the classic joint-venture firm: its profit fell eighty-eight percent in a single
year. And notice the tariffs: the EU put the highest duty on the most state-dependent firm. So
the market and the regulator point the same way. Next, the exact question we ask and the three
models we compare.""")


def s03_question(prs):
    s = blank(prs)
    kicker_headline(s, "Introduction · Research design",
                    "One question, three governance archetypes, four analytical dimensions")
    rect(s, MARGIN, Inches(1.78), Inches(12.1), Inches(0.9), NAVY)
    text(s, MARGIN + Inches(0.2), Inches(1.78), Inches(11.7), Inches(0.9),
         "“Which governance and financing structure is most likely to create long-term "
         "shareholder value in a technology-intensive, high-growth market?”",
         16, WHITE, font=HEAD, bold=True, italic=True, anchor=MSO_ANCHOR.MIDDLE, ls=1.05)
    cases = [("SAIC Motor", "Joint-venture model\nState-owned · SASAC ~73%", SAIC),
             ("BYD", "Vertical integration\nFounder-led · Wang ~17%", BYD_C),
             ("Geely", "M&A-driven\nFounder-led · Li ~41% attributable", GEELY)]
    x = MARGIN
    for title, desc, c in cases:
        rect(s, x, Inches(2.95), Inches(3.93), Inches(1.35), CARD, line_color=LINE)
        rect(s, x, Inches(2.95), Inches(3.93), Inches(0.09), c)
        text(s, x + Inches(0.18), Inches(3.13), Inches(3.6), Inches(0.4), title, 16, c, bold=True)
        text(s, x + Inches(0.18), Inches(3.58), Inches(3.6), Inches(0.7), desc, 13, INK, ls=1.05)
        x = Emu(x + Inches(4.08))
    text(s, MARGIN, Inches(4.6), Inches(12.1), Inches(0.3),
         "FOUR ANALYTICAL DIMENSIONS — APPLIED TO EACH CASE", 12.5, BLUE, bold=True)
    dims = ["1  Governance", "2  Financing", "3  Capability Building", "4  Long-term Viability"]
    x = MARGIN
    for i, d in enumerate(dims):
        rect(s, x, Inches(5.0), Inches(2.8), Inches(0.7), NAVY2)
        text(s, x, Inches(5.0), Inches(2.8), Inches(0.7), d, 13.5, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            text(s, Emu(x + Inches(2.8)), Inches(5.0), Inches(0.2), Inches(0.7), "›", 18,
                 GREY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = Emu(x + Inches(3.0))
    text(s, MARGIN, Inches(6.05), Inches(12.1), Inches(0.6),
         "Central thesis: governance structure, more than initial technological endowment, "
         "determines long-term value creation.", 14, NAVY, bold=True, ls=1.05)
    footer(s, 3)
    notes(s, """Our question is short to say but hard to answer: which structure builds the most
shareholder value over time, in a market that is both high-growth and technology-intensive. We
study it through three firms that picked three different paths. SAIC is state-owned and grew
through joint ventures. BYD is founder-led and built its technology in-house. Geely is also
founder-led but grew by buying companies abroad, like Volvo. To keep the comparison fair, we look
at all three through the same four lenses: governance, financing, capability building, and
long-term survival. The same four questions, asked of every firm. Our thesis, in one line: the
structure of control matters more than the technology a firm starts with. Now the theory that
stands behind these lenses.""")


def s04_theory(prs):
    s = blank(prs)
    kicker_headline(s, "Theory & literature",
                    "Six established lenses frame governance as the driver of value")
    th = [("Control rights drive allocation", "Shleifer & Vishny (1994)",
           "State control steers firms toward political goals, away from value.", RED),
          ("State vs. private ownership", "Megginson & Netter (2001)",
           "Privatised firms outperform; mixed-ownership firms do not catch up.", RED),
          ("Agency & founder alignment", "Jensen & Meckling (1976)",
           "Concentrated founder ownership removes the principal–agent conflict.", BLUE),
          ("Quid pro quo transfer", "Holmes et al. (2015) · Fan et al. (2007)",
           "Tech-for-access; politically connected firms lag ~18% post-IPO.", RED),
          ("Springboard M&A", "Luo & Tung (2007)",
           "Emerging-market firms acquire abroad to leap latecomer disadvantages.", TEAL),
          ("Dynamic capabilities", "Teece et al. (1997)",
           "Path-dependent tacit knowledge is the source of non-imitable advantage.", BLUE)]
    cw, ch, gx, gy = Inches(3.93), Inches(2.05), Inches(0.165), Inches(0.2)
    for i, (t, cite, body, c) in enumerate(th):
        col, row = i % 3, i // 3
        x = Emu(MARGIN + col * (cw + gx)); y = Emu(Inches(1.9) + row * (ch + gy))
        rect(s, x, y, cw, ch, CARD, line_color=LINE)
        rect(s, x, y, Inches(0.09), ch, c)
        text(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.16)), cw - Inches(0.36), Inches(0.5),
             t, 14.5, NAVY, bold=True, ls=0.95)
        text(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.78)), cw - Inches(0.36), Inches(0.3),
             cite, 11.5, c, bold=True, italic=True)
        text(s, Emu(x + Inches(0.22)), Emu(y + Inches(1.12)), cw - Inches(0.36), Inches(0.85),
             body, 12.5, INK, ls=1.0)
    footer(s, 4)
    notes(s, """Our argument rests on six well-known ideas, and they all point the same way. Shleifer
and Vishny show that whoever holds control rights decides where the money goes — and in a state
firm that means political goals, not shareholder value. Megginson and Netter add the evidence:
private firms beat state and mixed-ownership firms. Jensen and Meckling explain the opposite case
— a strong founder-owner removes the conflict between owners and managers. Holmes and Fan describe
the joint-venture trade and its hidden cost: connected firms lag about eighteen percent after
listing. Luo and Tung explain why an emerging-market firm buys abroad — to leap over its late
start. And Teece explains why deep, in-house knowledge is so hard for rivals to copy. Together
they predict that concentrated, commercial control should adapt faster and create more value.
Next we turn these six lenses into one simple map.""")


def s05_framework(prs):
    s = blank(prs)
    kicker_headline(s, "Analytical framework",
                    "Three models tested against four dimensions — a 3 × 4 matrix")
    headers = ["", "Governance", "Financing", "Capability", "Viability"]
    rows = [["", "State / SASAC + Party", "JV dividends\n(capital-light)", "Mandatory transfer", "Most exposed"],
            ["", "Founder-led\n(Wang ~17%)", "Self-financing\n(operating cash)", "Organic build", "Strongest"],
            ["", "Founder-led\n(Li ~41%)", "Leveraged M&A", "Acquisition + co-dev", "Most asymmetric"]]
    col_ws = [Inches(1.7), Inches(2.85), Inches(2.55), Inches(2.55), Inches(2.45)]
    cc = {(0, 0): RED, (1, 0): BLUE, (2, 0): TEAL}
    top, rh = Inches(1.95), Inches(0.92)
    grid(s, top, MARGIN, col_ws, rh, headers, rows, fs=13.5, hs=13.5, cell_colors=cc)
    for ri, (lbl, c) in enumerate([("SAIC", RED), ("BYD", BLUE), ("Geely", TEAL)]):
        y = Emu(top + rh * (ri + 1))
        text(s, MARGIN, y, col_ws[0], rh, lbl, 14, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, MARGIN, Inches(5.95), Inches(12.1), Inches(0.6),
         "Reading the matrix: the firm that is least adaptable in governance is also the most "
         "exposed in viability — the columns are linked, not independent.", 14, NAVY, bold=True, ls=1.05)
    source(s, "Anchors: Shleifer & Vishny (1994) · Megginson & Netter (2001) · Teece et al. (1997) · Luo & Tung (2007) · Fan et al. (2007)")
    footer(s, 5)
    notes(s, """This matrix is the map for the rest of the talk, so it is worth a moment. The three
rows are the three firms. The four columns are our four lenses. If you read across a row, you see
one firm's whole profile. If you read down a column, you compare the three firms on one topic.
The key message is that the columns are not separate. Look at SAIC: state-led governance,
capital-light joint-venture financing, technology that was handed to it, and the weakest position
in the EV shift. Those four cells belong together — the governance choice in column one shapes the
outcome in column four. BYD shows the opposite pattern, all the way across. Over the next slides
we go column by column and show the evidence behind each cell. We begin with governance.""")


def s06_governance(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 1 · Governance",
                    "Who controls the firm — and what they optimise — diverges sharply")
    card(s, MARGIN, Inches(1.85), Inches(3.93), Inches(4.15), SAIC, "SAIC — dual state principal", [
        "SASAC owns ~73%; the Party mandate is written into the charter",
        "Optimises beyond shareholder value: jobs, policy, prestige",
        "Connected firms lag ~18% post-IPO (Fan et al., 2007)",
        "Cannot exit unprofitable JVs without state sign-off",
        "Slow capital reallocation → strategic rigidity"], bsize=12.5)
    card(s, Emu(MARGIN + Inches(4.08)), Inches(1.85), Inches(3.93), Inches(4.15), BYD_C, "BYD — founder control", [
        "Wang Chuanfu ~17%; founding team > 1/3 of votes",
        "No JV dependency; one clear commercial objective",
        "Founder control aids breakthrough innovation (Zhang et al., 2025)",
        "Exited combustion engines in 2022 — decided alone",
        "Fastest decision cycle of the three"], bsize=12.5)
    card(s, Emu(MARGIN + Inches(8.16)), Inches(1.85), Inches(3.93), Inches(4.15), GEELY, "Geely — founder apex", [
        "Li Shufu ~41% attributable via ZGH (64.4% combined)",
        "No state committee or JV-partner consent at the top",
        "Agile in M&A, but coordination cost across brands",
        "Volvo / Polestar / Lotus add portfolio complexity",
        "Loss-making subsidiaries = contingent liability"], bsize=12.5)
    source(s, "Theory: Shleifer & Vishny (1994) · Fan et al. (2007) · Jensen & Meckling (1976) · Megginson & Netter (2001)")
    footer(s, 6)
    notes(s, """Governance is really two questions: who holds control, and what they try to maximise.
SAIC sits under the state. The asset commission owns about seventy-three percent, and the Party
role is written into the company charter, so its goals go beyond profit — jobs, industrial policy,
and prestige all count. The cost is speed: it cannot leave a failing joint venture without state
approval. BYD is the opposite. The founder and his team hold more than a third of the votes, there
is one commercial goal, and that is why he could end petrol-car production in a single year, in
2022, without asking anyone. Geely is also founder-led and quick on deals, but it pays a price for
running many brands — Volvo, Polestar, Lotus — and some of them lose money. So control runs on a
scale from political to commercial. The next slide shows that the EU put an actual price on this
difference.""")


def s07_gov_synth(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 1 · Governance synthesis",
                    "The EU tariff gap prices governance risk — an outside validation")
    tar = [("35.3%", "SAIC", "State-owned;\nhighest subsidy finding", RED),
           ("17.0%", "BYD", "Founder-led;\nlowest state finding", BLUE),
           ("18.8%", "Geely", "Private;\nmoderate state finding", TEAL)]
    x = Emu(MARGIN + Inches(0.35))
    for rate, name, note_, c in tar:
        rect(s, x, Inches(1.9), Inches(3.6), Inches(2.2), CARD, line_color=LINE)
        rect(s, x, Inches(1.9), Inches(3.6), Inches(0.09), c)
        text(s, x, Inches(2.12), Inches(3.6), Inches(0.95), rate, 44, c, font=HEAD,
             bold=True, align=PP_ALIGN.CENTER)
        text(s, x, Inches(3.05), Inches(3.6), Inches(0.4), name, 16, NAVY, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, Emu(x + Inches(0.12)), Inches(3.45), Inches(3.36), Inches(0.55), note_,
             12, GREY, italic=True, align=PP_ALIGN.CENTER, ls=1.0)
        x = Emu(x + Inches(3.85))
    bullets(s, MARGIN, Inches(4.45), Inches(12.1), Inches(1.9), [
        "An independent regulator — the European Commission — priced SAIC's state dependence at an 18.3-point premium over BYD, partly on “facts available” after the group withheld financing data.",
        "Adaptability rises as control concentrates: JV boards (bilateral consent) → state oversight → founder control.",
        "Core message: governance is a first-order determinant of competitive position, not an administrative footnote.",
    ], 15, gap=9)
    source(s, "European Commission (2024), Implementing Regulation (EU) 2024/2754 · Megginson & Netter (2001)")
    footer(s, 7)
    notes(s, """Here is outside proof that governance is not just an internal matter. The European
Commission studied how much state support each firm receives and set import duties accordingly:
35.3 percent on SAIC, 17 on BYD, 18.8 on Geely. That gap of about eighteen points between SAIC
and BYD is not about how good the cars are — it is about ownership. The more a firm leans on the
state, the higher the duty. In SAIC's case the Commission even had to use, in its own words,
“facts available,” because the group held back some of its financing information. So a neutral,
third-party regulator reached the same conclusion we do from the inside: control structure drives
competitive position. For a firm that wants to grow through exports, that tariff is a direct,
permanent cost of the governance model it chose. Now we follow the money — how each firm pays for
growth, and who carries the risk.""")


def s08_financing(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 2 · Financing",
                    "Capital structure decides who survives the downturn")
    cols = [("Capital-light JV", SAIC,
             ["Risk externalised to 50/50 JV partners",
              "SAIC-GM sales −56.5% triggered impairments",
              "Net profit −88% (2024); ROE just 0.58%",
              "Bound to absorb losses it cannot restructure",
              "→ a capital-light model became a capital trap"]),
            ("Self-financing", BYD_C,
             ["Funded from operating cash, not lenders",
              "R&D (RMB 54.2 bn) exceeded net profit (40.3 bn)",
              "USD 5.6 bn raise with sovereign wealth funds",
              "Net-cash position = resilience in a downturn",
              "→ resilient, but growth capped by internal cash"]),
            ("Leveraged M&A", GEELY,
             ["Debt-financed acquisitions buy capability fast",
              "Volvo (2010) for RMB 10.2 bn; then Lynk & Co",
              "Gearing rose 8.8% → 19.8% (FY2025)",
              "Operating cash flow RMB 47.3 bn gives headroom",
              "→ speed, at the cost of balance-sheet exposure"])]
    x = MARGIN
    for title, c, items in cols:
        card(s, x, Inches(1.85), Inches(3.93), Inches(3.7), c, title, items, bsize=12.5)
        x = Emu(x + Inches(4.08))
    rect(s, MARGIN, Inches(5.78), Inches(12.1), Inches(0.85), NAVY)
    text(s, MARGIN + Inches(0.2), Inches(5.78), Inches(11.7), Inches(0.85),
         "Risk-sharing captures upside in growth but exposes the firm in a downturn; self-financing "
         "trades growth speed for resilience; leverage trades balance-sheet safety for strategic speed.",
         13.5, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, ls=1.05)
    source(s, "SAIC AR 2024/2025 · BYD AR 2024/2025 · Geely AR 2024/2025 · Balcet et al. (2012)")
    footer(s, 8)
    notes(s, """Each firm pays for growth in a different way, and the difference only shows up when
the market turns down. SAIC chose a capital-light model: it shared the cost and the risk fifty-fifty
with its foreign partners. That felt safe while sales were rising. But when its main joint venture,
SAIC-GM, lost more than half its sales, SAIC was still contractually bound to absorb the losses,
without the control to fix the business. Profit fell eighty-eight percent and return on equity
dropped to barely half a percent — the capital-light model had become a capital trap. BYD pays for
itself out of its own cash, and it even spends more on research than it earns in profit, so it does
not depend on banks; the trade-off is that its growth is limited by the cash it generates. Geely
borrows to buy capability quickly — that is how it got Volvo — but its debt is rising. So the
choice is really resilience versus speed. The next slide puts the hard numbers side by side.""")


def s09_table(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 2 · Financial profile",
                    "The numbers confirm it: BYD self-funds, SAIC leans on eroding JV rents")
    headers = ["Metric (FY24 / FY25)", "SAIC", "BYD", "Geely"]
    rows = [["Revenue (RMB bn)", "627.6 / 656.2", "777.1 / 804.0", "240.2 / 345.2"],
            ["Gross margin", "9.4% / 10.1%", "19.4% / 17.7%", "15.9% / 16.6%"],
            ["Net profit (RMB bn)", "1.7 / 10.1", "40.3 / 32.6", "16.8 / 16.9"],
            ["Return on equity", "0.58% / 3.43%", "21.7% / 13.2%", "19.2% / 18.2%"],
            ["R&D expense (RMB bn)", "35.2 / 33.6", "54.2 / 63.4", "10.4 / 17.6"],
            ["Operating cash flow", "69.3 / 34.3", "133.5 / 59.1", "26.5 / 47.3"],
            ["EU anti-subsidy tariff", "35.3%", "17.0%", "18.8%"]]
    col_ws = [Inches(3.6), Inches(2.83), Inches(2.83), Inches(2.84)]
    grid(s, Inches(1.85), MARGIN, col_ws, Inches(0.56), headers, rows, fs=13, hs=13.5)
    for i, c in enumerate([SAIC, BYD_C, GEELY]):
        xx = Emu(MARGIN + col_ws[0] + col_ws[1] * i)
        rect(s, xx, Inches(1.85), col_ws[1 + i], Inches(0.56), c)
        text(s, xx, Inches(1.85), col_ws[1 + i], Inches(0.56), ["SAIC", "BYD", "Geely"][i],
             13.5, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, MARGIN, Inches(6.35), Inches(12.1), Inches(0.45),
         "BYD's R&D exceeds its net profit; SAIC's headline 2025 recovery is mostly the absence of "
         "2024 write-downs, not a structural turnaround.", 13, NAVY, bold=True, ls=1.0)
    source(s, "Table 1 (paper, §3.4) · Annual Reports SAIC/BYD/Geely 2024 & 2025 · EU Reg. 2024/2754")
    footer(s, 9)
    notes(s, """This is the financial picture in one table, taken straight from our paper. Look at a few
rows. Return on equity: BYD and Geely are near twenty percent, while SAIC is below one — that single
line captures the whole story. Research spending: BYD spends more on research and development than it
earns in profit, which is very unusual and shows how much it reinvests. Operating cash flow: BYD
generates far more cash from its operations than the other two. One honest note about SAIC: its
profit looks like it recovered in 2025, from 1.7 to 10.1 billion, but that is mostly because the big
write-downs of 2024 did not repeat — the underlying business has stabilised, not turned around. And
the last row, the EU tariff, is the outside marker of the same governance gap we saw earlier. With
the money side clear, we move to the third lens: how each firm builds its technology.""")


def s10_capability(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 3 · Capability building",
                    "How a firm learns determines what it can do next")
    cols = [("SAIC — mandatory transfer", SAIC,
             ["Tech-for-access via state-mandated JVs",
              "+8.3% quality, but rents killed own R&D (Bai 2025; Howell 2018)",
              "No proprietary EV platform when the cycle turned"]),
            ("BYD — organic build", BYD_C,
             ["Built in-house: battery → power electronics → software",
              "Path-dependent, non-imitable knowledge (Teece, 1997)",
              "Slowest to build, but the deepest moat"]),
            ("Geely — acquisition + co-dev", GEELY,
             ["Volvo (2010): engineering depth and brand at once",
              "Real co-development, not copying (CMA platform)",
              "Patent convergence confirms co-creation (Konda, 2022)"])]
    x = MARGIN
    for title, c, items in cols:
        card(s, x, Inches(2.15), Inches(3.93), Inches(2.7), c, title, items, bsize=13.5)
        x = Emu(x + Inches(4.08))
    text(s, MARGIN, Inches(5.35), Inches(12.1), Inches(0.6),
         "The mode of learning sets durability: handed-over technology fades when the cycle turns; "
         "built or co-developed knowledge lasts.", 14, NAVY, bold=True, ls=1.05)
    source(s, "Bai et al. (2025) · Howell (2018) · Teece et al. (1997) · Whitfield & Wuttke (2026) · Konda et al. (2022)")
    footer(s, 10)
    notes(s, """How a firm learns shapes what it can build later, and our three firms learned in three
different ways. SAIC took technology through its joint ventures. That worked: a careful study by Bai
and co-authors finds it raised product quality by about eight percent. But there was a hidden cost.
The easy profit from the foreign brands removed the pressure to invent its own technology — a former
minister called these rents “like opium.” So when the EV era arrived, SAIC had good cars but no
battery, power-electronics, or software platform of its own. BYD took the opposite road. It built
everything itself over decades, from battery cells up to software, often ahead of government policy.
That kind of deep, hands-on knowledge is exactly what Teece says rivals cannot copy. Geely found a
third way: it bought Volvo and turned the deal into real shared engineering, and patent data confirm
genuine co-creation, not one-way copying. So we have three speeds and three depths of learning. The
next slide shows the trade-off, and what the shallow path cost SAIC.""")


def s11_cap_synth(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 3 · Speed vs. depth",
                    "Capability building trades speed against depth against independence")
    rowsd = [(SAIC, "SAIC", "Fast, but shallow",
              "Fastest to quality (+8.3%) — but dependent, with no EV platform of its own."),
             (GEELY, "Geely", "Fast and deep",
              "Turned the Volvo purchase into real co-development (CMA platform; patent convergence)."),
             (BYD_C, "BYD", "Slow, but deepest",
              "30 years from cells to software — non-imitable, fully its own (Teece, 1997).")]
    y = Inches(2.0)
    for c, firm, verdict, body in rowsd:
        rect(s, MARGIN, y, Inches(12.1), Inches(0.95), CARD, line_color=LINE)
        rect(s, MARGIN, y, Inches(1.7), Inches(0.95), c)
        text(s, MARGIN, y, Inches(1.7), Inches(0.95), firm, 16, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, MARGIN + Inches(1.95), y, Inches(3.2), Inches(0.95), verdict, 15, c, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, MARGIN + Inches(5.2), y, Inches(6.7), Inches(0.95), body, 13.5, INK,
             anchor=MSO_ANCHOR.MIDDLE, ls=1.05)
        y = Emu(y + Inches(1.08))
    rect(s, MARGIN, Inches(5.4), Inches(12.1), Inches(1.0), NAVY)
    text(s, MARGIN + Inches(0.25), Inches(5.4), Inches(11.6), Inches(1.0),
         "The proof: SAIC-GM plants run at just 37% capacity and VW's MEB plant at 14% (2025) — "
         "dependency hardened into a capability deficit, while organic depth and co-development kept autonomy.",
         14, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, ls=1.1)
    source(s, "Bai et al. (2025) · Konda et al. (2022) · Teece et al. (1997) · SAIC Motor (2026)")
    footer(s, 11)
    notes(s, """This chart sums up learning as a trade-off between three things: speed, depth, and
independence. Moving right means faster; moving up means deeper and more independent. SAIC sits low
and to the right — it gained quality quickly, but the knowledge stayed shallow and dependent on
partners. BYD sits high but to the left — slow to build, yet deep and fully its own. Geely is in the
middle-right: fairly fast and fairly deep, because it turned a purchase into real shared engineering.
The cost of the shallow path shows in the numbers on the right. SAIC's joint-venture plants now run
at just 14 to 37 percent of capacity. That is what happens when dependency hardens into a real
capability deficit — there is no in-house platform to fall back on. So the way you build capability
is not just history; it decides who enters the EV era from strength. That sets up the fourth lens:
the EV transition itself.""")


def s12_viability(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 4 · Long-term viability",
                    "The EV transition tests which model was built to adapt")
    cols = [("SAIC — most exposed", SAIC,
             ["JV plants at 14–37% capacity",
              "But own-brand NEV +33%, with zero subsidies",
              "A race: eroding JV rents vs. thin margin"]),
            ("BYD — strongest position", BYD_C,
             ["No combustion engines since 2022; #1 NEV 4 years",
              ">1 m exports (+140%), 119 countries",
              "Self-funded; only margin under pressure"]),
            ("Geely — most asymmetric", GEELY,
             ["NEV 51.5% in H1 2025; H1 profit +102%",
              "Strong cash flow gives headroom",
              "But partner-dependent; Polestar/Lotus losses"])]
    x = MARGIN
    for title, c, items in cols:
        card(s, x, Inches(2.15), Inches(3.93), Inches(2.7), c, title, items, bsize=13.5)
        x = Emu(x + Inches(4.08))
    text(s, MARGIN, Inches(5.35), Inches(12.1), Inches(0.6),
         "Same technology shock, three very different outcomes — exactly what the governance and "
         "capability choices predicted.", 14, NAVY, bold=True, ls=1.05)
    source(s, "Annual Reports 2025 (FY2024) & 2026 (FY2025) · Geely interim results H1 2025 · European Commission (2024)")
    footer(s, 12)
    notes(s, """The EV transition works like an audit — it tests every earlier choice. SAIC is the most
exposed, with its joint-venture plants running far below capacity. But we want to be fair: its own
brand is making real progress. Electric sales grew thirty-three percent, and importantly with no
direct vehicle subsidies, so that growth stands on commercial footing. It even launched a new brand
with Huawei using a first-of-its-kind battery. So SAIC's future is a genuine race — falling
joint-venture income against a still-thin own-brand margin — and its state structure decides how fast
it can move. BYD is clearly the strongest: no petrol cars since 2022, four years as the world's top
EV maker, and exports up one hundred and forty percent. Its only real pressure is margin, from the
price war. Geely is the in-between case — very strong sales growth, but it leans on partners for
platforms and carries loss-making brands. Same shock, three very different results. Let us score them
side by side.""")


def s13_scorecard(prs):
    s = blank(prs)
    kicker_headline(s, "Synthesis · Structural audit",
                    "Across all four dimensions, private founder control outperforms")
    crit = ["Governance", "Financing", "Capability", "Viability", "Overall"]
    data = {"SAIC": [0, 0, 0, 0, 0], "BYD": [2, 2, 2, 2, 2], "Geely": [2, 1, 2, 1, 1]}
    cmap = {2: GREEN, 1: AMBER, 0: RED}; lab = {2: "Strong", 1: "Mixed", 0: "Weak"}
    compcol = {"SAIC": RED, "BYD": BLUE, "Geely": TEAL}
    top, labw, cw, rh = Inches(2.0), Inches(1.9), Inches(2.04), Inches(0.95)
    rect(s, MARGIN, top, labw, Inches(0.55), NAVY)
    x = Emu(MARGIN + labw)
    for c in crit:
        rect(s, x, top, cw, Inches(0.55), NAVY)
        text(s, Emu(x + Inches(0.05)), top, cw - Inches(0.1), Inches(0.55), c, 12, WHITE,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=0.9)
        x = Emu(x + cw)
    for ri, comp in enumerate(["SAIC", "BYD", "Geely"]):
        y = Emu(top + Inches(0.55) + ri * rh)
        rect(s, MARGIN, y, labw, rh, compcol[comp])
        text(s, MARGIN, y, labw, rh, comp, 16, WHITE, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        x = Emu(MARGIN + labw)
        for ci in range(len(crit)):
            rect(s, x, y, cw, rh, WHITE, line_color=LINE)
            v = data[comp][ci]; d = Inches(0.3)
            oval(s, Emu(x + cw / 2 - d / 2), Emu(y + Inches(0.16)), d, cmap[v])
            text(s, x, Emu(y + Inches(0.52)), cw, Inches(0.3), lab[v], 10.5, cmap[v],
                 bold=True, align=PP_ALIGN.CENTER)
            x = Emu(x + cw)
    text(s, MARGIN, Inches(6.05), Inches(12.1), Inches(0.7),
         "Architecture decides EV-transition speed: founder control + self-financing (BYD) turns the "
         "shock into leadership; state + JV dependence (SAIC) turns it into a structural discount.",
         13.5, NAVY, bold=True, ls=1.05)
    lx = W - Inches(4.6)
    for i, (v, lb) in enumerate([(2, "Strong"), (1, "Mixed"), (0, "Weak")]):
        oval(s, Emu(lx + Inches(1.5) * i), H - Inches(0.46), Inches(0.18), cmap[v])
        text(s, Emu(lx + Inches(1.5) * i + Inches(0.24)), H - Inches(0.5), Inches(1.2),
             Inches(0.3), lb, 10, GREY)
    footer(s, 13)
    notes(s, """This scorecard puts the whole comparison on one page. Green means strong, amber means
mixed, red means weak, and the four columns are our four lenses plus an overall column. The pattern
is hard to miss. BYD is green all the way across — strong governance, strong financing, deep
capability, and the best position for the EV shift. SAIC is red across the board, weakest on the
ability to adapt, because the state structure slows every decision. Geely sits in between: strong on
governance and capability, but only mixed on financing and viability because of its debt and its
loss-making brands. The deeper point is the one we keep returning to: the structure a firm chose
years ago decides how fast it can move now. Founder control plus self-financing turned the
technology shock into leadership for BYD; state ownership plus joint-venture dependence turned the
same shock into a lasting discount for SAIC. With that, we can answer our central question.""")


def s14_verdict(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, BLUE)
    text(s, MARGIN, Inches(0.5), Inches(12), Inches(0.3), "VERDICT", 14, TEAL, bold=True)
    text(s, MARGIN, Inches(0.85), Inches(12.1), Inches(1.0),
         "BYD's founder-led vertical integration is the model most likely\nto create long-term "
         "shareholder value", 26, WHITE, font=HEAD, bold=True, ls=1.0)
    metrics = [("21.7%", "highest ROE of the three\nfirms (FY2024)"),
               ("4 yrs", "consecutive global NEV\nmarket leader"),
               ("USD 5.6 bn", "2025 raise incl. sovereign\nwealth funds")]
    x = MARGIN
    for val, lbl in metrics:
        rect(s, x, Inches(2.5), Inches(3.93), Inches(1.5), NAVY2)
        text(s, x, Inches(2.65), Inches(3.93), Inches(0.7), val, 32, WHITE, font=HEAD,
             bold=True, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.15), Inches(3.42), Inches(3.63), Inches(0.5), lbl, 12,
             RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER, ls=1.0)
        x = Emu(x + Inches(4.08))
    rect(s, MARGIN, Inches(4.35), Inches(12.1), Inches(1.6), NAVY2)
    rect(s, MARGIN, Inches(4.35), Inches(0.09), Inches(1.6), TEAL)
    text(s, MARGIN + Inches(0.25), Inches(4.55), Inches(11.6), Inches(1.25),
         "Why BYD — and why not simply “vertical integration is best”: the model wins because it "
         "aligns concentrated, commercially-oriented control with a self-financing structure that "
         "stays resilient when markets contract — not because integration is inherently superior. "
         "The advantage is governance and capital discipline working together.",
         14.5, WHITE, ls=1.15)
    text(s, MARGIN, Inches(6.2), Inches(12.1), Inches(0.4),
         "Consistent with Allen et al. (2024): privately governed Chinese firms structurally "
         "outperform state-controlled peers. ROE compressed to 13.2% in FY2025 under price-war pressure.",
         11, GREY, italic=True, ls=1.0)
    footer(s, 13)
    notes(s, """Here is our answer to the central question. Of the three models, BYD is the one most
likely to create long-term shareholder value. The evidence is consistent: the highest return on
equity of the three, four straight years as the world's top EV maker, and a five-point-six-billion
dollar capital raise that sovereign wealth funds chose to join — a strong outside vote of confidence.
But please notice the careful wording in the box. We are not saying that building everything yourself
is automatically the best strategy. BYD wins because it combines two things at once: concentrated,
commercial founder control, and a self-financing balance sheet that stays strong when the market
turns down. Governance and capital discipline working together — that is the real driver. This also
matches the broader finding by Allen and co-authors that private Chinese firms outperform state ones.
And we stay honest: BYD's own return fell to thirteen percent in 2025 under the price war, so the
model is strong but not untouchable. From this single answer we can now draw three lessons that
travel beyond this industry.""")


def s15_lessons(prs):
    s = blank(prs)
    kicker_headline(s, "Generalizable lessons",
                    "Three transferable lessons on governance and value in high-growth markets")
    lessons = [("1", "Capability mode determines durability across technology cycles",
                "Mandatory transfer buys short-run quality but displaces the autonomous innovation the next cycle needs; organic and voluntary modes preserve it.", BLUE),
               ("2", "Governance is a first-order determinant of strategic adaptability",
                "Control concentrated with commercial principals adapts faster. When change outpaces the decision cycle, governance becomes the binding constraint.", RED),
               ("3", "Capital structure determines resilience when markets contract",
                "Risk-sharing captures upside in growth but exposes the firm in downturns; self-financing trades speed for resilience, leverage trades safety for speed.", TEAL)]
    y = Inches(1.95)
    for num, title, body, c in lessons:
        rect(s, MARGIN, y, Inches(12.1), Inches(1.3), CARD, line_color=LINE)
        rect(s, MARGIN, y, Inches(0.12), Inches(1.3), c)
        text(s, MARGIN + Inches(0.3), y, Inches(1.0), Inches(1.3), num, 44, c, font=HEAD,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
        text(s, MARGIN + Inches(1.4), Emu(y + Inches(0.18)), Inches(10.4), Inches(0.5),
             title, 16.5, NAVY, bold=True, ls=1.0)
        text(s, MARGIN + Inches(1.4), Emu(y + Inches(0.62)), Inches(10.4), Inches(0.6),
             body, 13, INK, ls=1.05)
        y = Emu(y + Inches(1.5))
    footer(s, 14)
    notes(s, """From these three cases we draw three lessons that reach beyond the car industry, into any
fast-moving, uncertain market. First: the way you build capability decides how long it lasts. Buying
or borrowing technology can give quick results, but it can also hollow out a firm so it has nothing of
its own when the technology changes — that was SAIC. Building or genuinely co-developing keeps that
inner strength. Second: governance is a first-order driver of how fast a firm can adapt. When control
sits with commercial owners, the firm moves quickly; when it is split across political and partner
interests, it stalls — and when technology moves faster than your decisions, governance becomes the
thing that holds you back. Third: capital structure decides who survives a downturn. Sharing risk
feels safe in good times but traps you in bad times; financing yourself trades some speed for
resilience; heavy borrowing trades safety for speed. Three lessons, one theme — structure decides.
Let me now bring it all together in the conclusion.""")


def s16_conclusion(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, BLUE)
    text(s, MARGIN, Inches(0.5), Inches(12), Inches(0.3),
         "CONCLUSION", 14, TEAL, bold=True)
    text(s, MARGIN, Inches(0.95), Inches(12.1), Inches(0.6),
         "Structure, more than technology, decided the outcome", 28, WHITE, font=HEAD, bold=True)
    qx = MARGIN
    rect(s, qx, Inches(1.95), Inches(7.3), Inches(4.0), NAVY2)
    rect(s, qx, Inches(1.95), Inches(7.3), Inches(0.09), TEAL)
    text(s, qx + Inches(0.3), Inches(2.15), Inches(3.8), Inches(0.8), "“", 60, TEAL,
         font=HEAD, bold=True)
    text(s, qx + Inches(0.3), Inches(2.95), Inches(6.7), Inches(2.8),
         "When technology moves faster than firms can decide, governance is no longer a footnote "
         "to strategy — it is strategy.", 22, WHITE, font=HEAD, bold=True, italic=True, ls=1.12)
    bx = Emu(qx + Inches(7.7))
    text(s, bx, Inches(2.1), Inches(4.4), Inches(0.4), "What the evidence shows", 14, TEAL, bold=True)
    bullets(s, bx, Inches(2.6), Inches(4.4), Inches(3.3), [
        "The JV model fell not to a better technology, but to a better way of allocating control and capital.",
        "BYD — founder control + self-financing — was the most adaptive combination in the transition.",
        "SAIC's state model optimises beyond shareholder value: a permanent structural discount.",
        "2020–2025 is a natural experiment: governance, more than endowment, decides who wins.",
    ], 13, color=RGBColor(0xE2, 0xE8, 0xF0), gap=10, mcolor=TEAL)
    footer(s, 15)
    notes(s, """Let me bring the whole argument to a single point. Across all four lenses, the evidence
converges on one finding: governance structure, more than the technology a firm starts with, decides
long-term value creation. The joint-venture model that organised Chinese carmaking for thirty years
did not lose to a better battery — it lost to a better way of allocating control and capital. BYD's
combination of founder control and self-financing proved the most adaptive; Geely adapted too, but
at rising balance-sheet cost; and SAIC's state model, by design, serves goals beyond shareholder
value, which shows up as a permanent discount. The years 2020 to 2025 gave us something close to a
natural experiment — one technology shock, three governance structures, three clearly different
outcomes. So our closing line is this: when technology moves faster than firms can decide, governance
is no longer a footnote to strategy — it is strategy. Before questions, we want to be transparent
about the study's limits and about how we used AI.""")


def s17_limits(prs):
    s = blank(prs)
    kicker_headline(s, "Limitations & future research",
                    "Scope conditions and where the analysis could be extended")
    rect(s, MARGIN, Inches(1.95), Inches(5.95), Inches(3.85), CARD, line_color=LINE)
    rect(s, MARGIN, Inches(1.95), Inches(5.95), Inches(0.09), AMBER)
    text(s, MARGIN + Inches(0.2), Inches(2.15), Inches(5.55), Inches(0.4), "Limitations", 16,
         AMBER, bold=True)
    bullets(s, MARGIN + Inches(0.2), Inches(2.65), Inches(5.55), Inches(3.0), [
        "Three firms, one industry, one country — limited statistical generalisability.",
        "Short post-EV window; the transition is still unfolding.",
        "Reliance on reported figures; currency conventions (RMB vs. USD) differ across reports.",
        "NIO's capital-market model is noted but outside this team's scope.",
    ], 13, gap=9, mcolor=AMBER)
    lx = Emu(MARGIN + Inches(6.15))
    rect(s, lx, Inches(1.95), Inches(5.95), Inches(3.85), CARD, line_color=LINE)
    rect(s, lx, Inches(1.95), Inches(5.95), Inches(0.09), BLUE)
    text(s, Emu(lx + Inches(0.2)), Inches(2.15), Inches(5.55), Inches(0.4), "Future research",
         16, BLUE, bold=True)
    bullets(s, Emu(lx + Inches(0.2)), Inches(2.65), Inches(5.55), Inches(3.0), [
        "Add a fourth archetype (NIO, capital-market model) for a fuller typology.",
        "Track outcomes over a longer EV cycle to test durability of the ranking.",
        "Build a formal governance index and test it across more emerging markets.",
        "Examine how EU tariffs reshape each model's export strategy.",
    ], 13, gap=9, mcolor=BLUE)
    footer(s, 16)
    notes(s, """A short, honest word on limits before the AI part, since questions usually follow. We
study three firms in one industry and one country, so we cannot claim a statistical law — these are
deep case studies, not a large sample. The EV transition is also still young, so today's ranking
could shift as the story plays out. We rely on the figures companies report, and those reports mix
currencies, which we had to reconcile carefully. And we left NIO out — it represents a fourth,
capital-market model that would round out the picture but was beyond our scope. That last point is
also where future work could go: add NIO for a fuller typology, follow these firms over a longer
cycle, build a proper numerical governance score, and study how the EU tariffs reshape each firm's
export plans. Now, as the seminar asks, we reflect openly on how we used AI in this project.""")


def s18_ai_arch(prs):
    s = blank(prs)
    kicker_headline(s, "AI annex · Architecture & workflow",
                    "AI as research infrastructure — a layered agent system, not a chatbot")
    # Compact models + agents line (de-emphasised)
    text(s, MARGIN, Inches(1.8), Inches(12.1), Inches(0.3),
         "Platform: Claude Code (Anthropic) — claude-sonnet-4-6 (primary) · claude-opus-4 (hard analysis)",
         12, GREY, bold=True)
    text(s, MARGIN, Inches(2.18), Inches(12.1), Inches(0.3),
         "5 agents: wiki-ingest · wiki-query · wiki-lint · verify-claims · find-references",
         12, GREY, bold=True)
    # Prominent shared-repo folder tree (the core of the workflow)
    rect(s, MARGIN, Inches(2.75), Inches(6.5), Inches(3.55), CARD, line_color=LINE)
    rect(s, MARGIN, Inches(2.75), Inches(6.5), Inches(0.09), BLUE)
    text(s, MARGIN + Inches(0.25), Inches(2.95), Inches(6.0), Inches(0.4),
         "One shared git repository", 17, NAVY, bold=True)
    tree = [("research/wiki/", NAVY, 0),
            ("├─ saic-motor/", SAIC, 1),
            ("├─ byd/", BYD_C, 1),
            ("├─ geely/", GEELY, 1),
            ("└─ shared/   — EU regulation · IEA data", GREY, 1)]
    yy = Inches(3.5)
    for txt_, col, lvl in tree:
        text(s, MARGIN + Inches(0.4) + Inches(0.3) * lvl, yy, Inches(5.8), Inches(0.4),
             txt_, 15, col, bold=(lvl == 0), font="Consolas")
        yy = Emu(yy + Inches(0.46))
    # Right: why it matters
    rx = Emu(MARGIN + Inches(6.9))
    text(s, rx, Inches(2.95), Inches(5.4), Inches(0.4), "Why it mattered", 17, BLUE, bold=True)
    bullets(s, rx, Inches(3.5), Inches(5.4), Inches(2.8), [
        "One folder per case + a shared layer for common data",
        "Rule: git pull before each session, push after",
        "The tool read all three case wikis at once",
        "Direct cross-case comparison — no manual hand-offs",
    ], 14, gap=10)
    footer(s, 17)
    notes(s, """Now the AI reflection, which the seminar asks us to include. The first thing to say is
that we did not use a simple chatbot. We used Claude Code as a layered system. A main model did most
of the work — reading sources, drafting, checking citations — and a stronger model was brought in for
the hard analytical passes, like comparing the three cases and testing the logic of our argument.
Around these models we ran five small, specialised agents, each with one job: one to ingest sources,
one to search them, one to health-check the wiki for broken links and contradictions, one to trace
each claim back to its source, and one to find supporting literature. Everything lived in a single
shared git repository, with one folder per company and a shared layer for common data. Our simple
rule kept us in sync: pull before each session, push after. Because the tool could read all three
wikis at once, we compared the cases directly, without emailing files around. Next, what this setup
actually added to the quality of the work.""")


def s19_ai_value(prs):
    s = blank(prs)
    kicker_headline(s, "AI annex · Value & quality control",
                    "Where AI added the most value")
    # Positives first and prominent
    text(s, MARGIN, Inches(1.8), Inches(12.1), Inches(0.3), "WHAT WORKED IN PRACTICE", 13,
         BLUE, bold=True)
    bullets(s, MARGIN, Inches(2.15), Inches(12.1), Inches(2.0), [
        "Compressed research time: 30+ papers and reports ingested, indexed and made searchable.",
        "Source discipline: every source rated by quality (Q1–Q5) × relevance (R1–R4), arguments on Q1/Q2 only.",
        "Sharpest prompt: “review this as a Dozent looking for overclaims, logical gaps and APA errors.”",
        "Fact-checking: isolate one claim and make it quote the exact source passage before use.",
    ], 14, gap=9)
    # Errors as a smaller supporting row at the bottom
    text(s, MARGIN, Inches(4.55), Inches(12.1), Inches(0.3),
         "AND IT CAUGHT THREE DATA ERRORS BEFORE SUBMISSION", 12, GREEN, bold=True)
    errs = [("EU tariff", "Provisional 36.3% → final 35.3% (SAIC)."),
            ("Currency unit", "A net-profit figure inflated 10× in a source."),
            ("Share count", "A BYD EPS built on the wrong share count.")]
    x = MARGIN
    for i, (t, d) in enumerate(errs):
        rect(s, x, Inches(4.9), Inches(3.93), Inches(1.15), CARD, line_color=LINE)
        rect(s, x, Inches(4.9), Inches(0.09), Inches(1.15), GREEN)
        text(s, x + Inches(0.18), Inches(5.02), Inches(3.6), Inches(0.35), f"{i+1} · {t}",
             13, GREEN, bold=True)
        text(s, x + Inches(0.18), Inches(5.4), Inches(3.6), Inches(0.6), d, 12, INK, ls=1.0)
        x = Emu(x + Inches(4.08))
    source(s, "All three corrections traced back to the PwC- / Ernst & Young-audited annual reports.")
    footer(s, 18)
    notes(s, """The single most valuable thing the tool did was catch mistakes before they reached the
paper. It found three real data errors. First, an EU tariff we would have quoted as final was actually
the provisional figure — the correct one is 35.3 percent. Second, a net-profit number had been
inflated about ten times by a currency-unit mix-up in a secondary source. Third, an earnings-per-share
figure for BYD was built on the wrong share count. In every case we traced the number back to the
audited annual report and corrected it. On method, two habits made the difference. Telling the model
to act like a strict professor hunting for overclaims and citation errors gave much sharper feedback
than a friendly review. And forcing it to quote the exact source passage for each fact is what caught
those errors. We also rated every source by quality and only built arguments on the strongest ones.
Finally, the honest limits and the cost.""")


def s20_ai_limits(prs):
    s = blank(prs)
    kicker_headline(s, "AI annex · Limits, cost & judgment",
                    "Useful but bounded — and the thinking stayed human")
    rect(s, MARGIN, Inches(1.95), Inches(5.95), Inches(3.55), CARD, line_color=LINE)
    rect(s, MARGIN, Inches(1.95), Inches(5.95), Inches(0.09), AMBER)
    text(s, MARGIN + Inches(0.2), Inches(2.15), Inches(5.55), Inches(0.4), "Real limitations", 16,
         AMBER, bold=True)
    bullets(s, MARGIN + Inches(0.2), Inches(2.65), Inches(5.55), Inches(2.7), [
        "Cannot download paywalled PDFs — it knows the papers' content, but key claims still needed manual reading and checking.",
        "Geely PDFs were not machine-readable → more manual web research.",
        "Token / context limits hit quickly on the stronger model in long sessions.",
        "Output was often verbose and needed active editing; some DOIs were wrong.",
    ], 12.5, gap=8, mcolor=AMBER)
    lx = Emu(MARGIN + Inches(6.15))
    rect(s, lx, Inches(1.95), Inches(5.95), Inches(3.55), CARD, line_color=LINE)
    rect(s, lx, Inches(1.95), Inches(5.95), Inches(0.09), BLUE)
    text(s, Emu(lx + Inches(0.2)), Inches(2.15), Inches(5.55), Inches(0.4), "Lessons & cost", 16,
         BLUE, bold=True)
    bullets(s, Emu(lx + Inches(0.2)), Inches(2.65), Inches(5.55), Inches(2.7), [
        "A learning curve: several hours to use the agent system productively, not as a chat tool.",
        "AI output always needs critical review — early drafts contained overclaims.",
        "Cost: Claude Pro at ~€21/month per member; no per-query charge, no paid data services.",
        "Framing that worked: a source-quality enforcer and critical reviewer, not a writer.",
    ], 12.5, gap=8, mcolor=BLUE)
    rect(s, MARGIN, Inches(5.75), Inches(12.1), Inches(0.7), NAVY)
    text(s, MARGIN + Inches(0.2), Inches(5.75), Inches(11.7), Inches(0.7),
         "Bottom line: AI compressed research time and caught errors — but every judgment, argument "
         "and conclusion remained the authors' own.", 13.5, WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, ls=1.0)
    footer(s, 19)
    notes(s, """We also want to be clear about what the tool could not do. It cannot open paywalled
journals, so some reading stayed manual. Geely's PDFs were not machine-readable, which is why that
case leaned more on web sources. On the stronger model we hit token and context limits fairly
quickly in long sessions, and the output was often too wordy, so it always needed editing — it even
produced some wrong DOIs that we fixed by hand. The biggest lesson was that this is not a chat box:
it took several hours to learn to use the agent system well, and every output needed a critical read.
The cost was modest, about twenty-one euros a month each, with no extra charges. The line we held
throughout is the one on the bar at the bottom: we used AI as research infrastructure and as a strict
reviewer, never as a ghostwriter. Every argument, interpretation, and conclusion is our own. With
that, let me show our sources, and then we are happy to take your questions.""")


def s21_refs(prs):
    s = blank(prs)
    kicker_headline(s, "References", "Key sources (APA 7)")
    refs_l = [
        "Allen, F., Qian, J., Shan, C., & Zhu, J. L. (2024). Dissecting the long-term performance of the Chinese stock market. Journal of Finance, 79(2).",
        "Bai, J., Barwick, P. J., Cao, S., & Li, S. (2025). Quid pro quo, knowledge spillovers, and industrial quality upgrading. American Economic Review, 115(11).",
        "Balcet, G., Wang, H., & Richet, X. (2012). Geely: A trajectory of catching up and asset-seeking growth. IJATM, 12(4).",
        "European Commission. (2024). Implementing Regulation (EU) 2024/2754. Official Journal of the EU.",
        "Fan, J. P. H., Wong, T. J., & Zhang, T. (2007). Politically connected CEOs and post-IPO performance. JFE, 84(2).",
        "Holmes, T. J., McGrattan, E. R., & Prescott, E. C. (2015). Quid pro quo: Technology capital transfers for market access in China. RES, 82(3).",
        "Howell, S. T. (2018). Joint ventures and technology adoption: A Chinese industrial policy that backfired. Research Policy, 47(8).",
        "Jensen, M. C., & Meckling, W. H. (1976). Theory of the firm. JFE, 3(4).",
        "Konda, P., Slepnikov, D., & Jin, J. (2022). From transaction to co-creation in Geely's acquisition of Volvo Cars. AJTI, 31(3).",
    ]
    refs_r = [
        "Luo, Y., & Tung, R. L. (2007). International expansion of emerging market enterprises: A springboard perspective. JIBS, 38(4).",
        "Megginson, W. L., & Netter, J. M. (2001). From state to market: A survey of empirical studies on privatization. JEL, 39(2).",
        "Shleifer, A., & Vishny, R. W. (1994). Politicians and firms. Quarterly Journal of Economics, 109(4).",
        "Teece, D. J., Pisano, G., & Shuen, A. (1997). Dynamic capabilities and strategic management. SMJ, 18(7).",
        "Whitfield, L., & Wuttke, T. (2026). China's technological catch-up and leapfrogging in electric vehicles. Progress in Economic Geography, 4.",
        "Zhang, Z.-Q., Tang, B.-J., Su, Z., & Zhang, Y. (2025). Founder control and breakthrough innovation. China Economic Review, 94.",
        "Zheng, Q., Noorderhaven, N., & Du, J. (2022). Making the unlikely marriage work. Journal of World Business, 57(2).",
        "International Energy Agency. (2026). Global EV Outlook 2026. IEA.",
        "Annual Reports: SAIC Motor (2025, 2026); BYD (2025, 2026); Geely Automobile Holdings (2025, 2026).",
    ]
    bullets(s, MARGIN, Inches(1.85), Inches(5.95), Inches(5.0), refs_l, 10.5, gap=7, marker="", ls=1.0)
    bullets(s, Emu(MARGIN + Inches(6.15)), Inches(1.85), Inches(5.95), Inches(5.0), refs_r, 10.5,
            gap=7, marker="", ls=1.0)
    footer(s, 20)
    notes(s, """These are our main sources, in APA style, split across the two columns. The theory comes
from Shleifer and Vishny, Jensen and Meckling, Megginson and Netter, Luo and Tung, and Teece. The
empirical evidence on China comes from Bai, Howell, Holmes, and Allen. The company figures come from
the audited annual reports of all three firms, plus the EU regulation and the IEA outlook. The full
reference list is in the paper. Thank you very much for your attention — we are now happy to take your
questions.""")


def main():
    out_dir = r"C:\Users\User\CHINA-seminar-paper\research\presentation"
    out = os.path.join(out_dir, "China-Strategies-Presentation-FINAL.pptx")
    prs = new_prs()
    for fn in [s01_title, s02_intro, s03_question, s04_theory, s05_framework, s06_governance,
               s07_gov_synth, s08_financing, s09_table, s10_capability, s11_cap_synth,
               s12_viability, s14_verdict, s15_lessons, s16_conclusion,
               s17_limits, s18_ai_arch, s19_ai_value, s20_ai_limits, s21_refs]:
        fn(prs)
    prs.save(out)
    print("Saved:", out, "| slides:", len(prs.slides))


if __name__ == "__main__":
    main()
