import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ─────────────────────────────────────────────────────────────────────
WORKSPACE = r"C:/Users/Batu/research-template/research/presentation/workspace"
OUT_FILE  = r"C:/Users/Batu/research-template/research/presentation/Slides/Titel.pptx"
GRAD_PNG  = os.path.join(WORKSPACE, "gradient_stripe.png")

# ── Colors ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x11, 0x13, 0x18)
BLUE     = RGBColor(0x00, 0xC2, 0xFF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xB0, 0xB4, 0xBE)
GOLD     = RGBColor(0xD4, 0xAF, 0x37)
RED      = RGBColor(0xCC, 0x00, 0x00)
DIM_GRAY = RGBColor(0x70, 0x74, 0x7E)

# ── Slide: 16:9 widescreen ────────────────────────────────────────────────────
W = 13.333
H = 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def rect(x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    return s

def text(txt, x, y, w, h,
         fn="Tahoma", fs=12, bold=False,
         color=WHITE, align=PP_ALIGN.LEFT, spc=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.name = fn; run.font.size = Pt(fs)
    run.font.bold = bold; run.font.color.rgb = color
    if spc is not None:
        run._r.get_or_add_rPr().set('spc', str(int(spc * 100)))
    return tb

# ── Background ────────────────────────────────────────────────────────────────
rect(0, 0, W, H, BG)

# ── Gradient stripe (PNG overlay, full-slide, transparent) ────────────────────
if os.path.exists(GRAD_PNG):
    slide.shapes.add_picture(GRAD_PNG, Inches(0), Inches(0), Inches(W), Inches(H))

# ── Red left accent bar ───────────────────────────────────────────────────────
rect(0, 0, 0.055, H, RED)

# ── Layout anchors ────────────────────────────────────────────────────────────
# Top label:   y = 0.38
# Title block: y = 1.10  (two lines at 44pt ≈ 1.47in each = ~1.65in block)
# Subtitle:    y = 3.10  (below title block, with comfortable whitespace)
# Separator:   y = 6.50  (near bottom)
# Footer:      y = 6.70

LABEL_Y  = 0.38
TITLE_Y  = 1.10
SUBT_Y   = 3.10
SEP_Y    = 6.50
FOOT_Y   = 6.65
LEFT_X   = 0.50

# ── Top-left label ────────────────────────────────────────────────────────────
text("SEMINAR  ·  CHINA",
     LEFT_X, LABEL_Y, 4.0, 0.35,
     fs=9, color=BLUE, spc=2.5)

# ── Main title ────────────────────────────────────────────────────────────────
tb = slide.shapes.add_textbox(
    Inches(LEFT_X), Inches(TITLE_Y), Inches(12.3), Inches(3.0))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
p.space_after = Pt(2)
r = p.add_run()
r.text = "Governance-Modelle in der"
r.font.name = "Arial"; r.font.size = Pt(44)
r.font.bold = True; r.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.LEFT
p2.space_after = Pt(0)
r2 = p2.add_run()
r2.text = "chinesischen Automobilindustrie"
r2.font.name = "Arial"; r2.font.size = Pt(44)
r2.font.bold = True; r2.font.color.rgb = WHITE

# ── Subtitle ──────────────────────────────────────────────────────────────────
text("BYD  ·  Geely  ·  SAIC Motor  —  Ein strategischer Vergleich",
     LEFT_X, SUBT_Y, 12.3, 0.55,
     fn="Tahoma", fs=15, color=GRAY)

# ── Separator line (electric blue, full-width minus margins) ─────────────────
s = slide.shapes.add_shape(1,
    Inches(LEFT_X), Inches(SEP_Y), Inches(W - LEFT_X - 0.15), Inches(0.025))
s.fill.solid(); s.fill.fore_color.rgb = BLUE
s.line.fill.background()

# ── Footer left ───────────────────────────────────────────────────────────────
text("Seminar: CHINA  |  Sommersemester 2026",
     LEFT_X, FOOT_Y, 8.0, 0.42,
     fs=10, color=DIM_GRAY)

# ── Footer right (gold author) ────────────────────────────────────────────────
text("Batughan Maus",
     9.5, FOOT_Y, 3.5, 0.42,
     fn="Tahoma", fs=11, bold=True,
     color=GOLD, align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT_FILE), exist_ok=True)
prs.save(OUT_FILE)
print(f"Saved: {os.path.abspath(OUT_FILE)}")
