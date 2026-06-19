"""
Quick thumbnail renderer for the Titel.pptx using python-pptx + Pillow.
Renders shapes by reading their positions/sizes/colors and drawing them.
Good enough for layout verification.
"""
import os
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

PPTX_PATH = r"C:/Users/Batu/research-template/research/presentation/Slides/Titel.pptx"
OUT_PATH  = r"C:/Users/Batu/research-template/research/presentation/workspace/thumbnail.png"
GRAD_PNG  = r"C:/Users/Batu/research-template/research/presentation/workspace/gradient_stripe.png"

# Render at 2x for better quality (1440 x 810)
SCALE = 2
W_PX, H_PX = 720 * SCALE, 405 * SCALE

prs = Presentation(PPTX_PATH)
slide = prs.slides[0]

slide_w = prs.slide_width   # EMU
slide_h = prs.slide_height  # EMU

def emu_to_px(emu, slide_dim, px_dim):
    return int(emu / slide_dim * px_dim)

img = Image.new("RGB", (W_PX, H_PX), (17, 19, 24))
draw = ImageDraw.Draw(img)

for shape in slide.shapes:
    left   = shape.left   or 0
    top    = shape.top    or 0
    width  = shape.width  or 0
    height = shape.height or 0

    x0 = emu_to_px(left,          slide_w, W_PX)
    y0 = emu_to_px(top,           slide_h, H_PX)
    x1 = emu_to_px(left + width,  slide_w, W_PX)
    y1 = emu_to_px(top  + height, slide_h, H_PX)

    # ── Pictures (gradient PNG overlay) ──────────────────────────────────────
    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        try:
            from io import BytesIO
            img_data = shape.image.blob
            pic = Image.open(BytesIO(img_data)).convert("RGBA")
            pic = pic.resize((x1 - x0, y1 - y0), Image.LANCZOS)
            # Composite over existing image
            base = img.crop((x0, y0, x1, y1)).convert("RGBA")
            comp = Image.alpha_composite(base, pic)
            img.paste(comp.convert("RGB"), (x0, y0))
        except Exception:
            pass
        continue

    # ── Filled shapes (rectangles) ────────────────────────────────────────────
    fill_color = None
    try:
        if shape.fill.type is not None:
            fc = shape.fill.fore_color
            if fc.type is not None:
                rgb = fc.rgb
                fill_color = (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass

    if fill_color:
        draw.rectangle([x0, y0, x1, y1], fill=fill_color)

    # ── Text boxes ────────────────────────────────────────────────────────────
    if shape.has_text_frame:
        tf = shape.text_frame
        cur_y = y0
        for para in tf.paragraphs:
            line_text = ""
            line_color = (255, 255, 255)
            line_size  = 12
            line_bold  = False
            for run in para.runs:
                line_text += run.text
                try:
                    rgb = run.font.color.rgb
                    line_color = (rgb.red, rgb.green, rgb.blue)
                except Exception:
                    pass
                try:
                    if run.font.size:
                        line_size = int(run.font.size.pt)
                    line_bold = run.font.bold or False
                except Exception:
                    pass
            if not line_text.strip():
                cur_y += int(line_size * SCALE * 0.5)
                continue

            # Scale font size
            fs_px = int(line_size * SCALE * 0.75)  # pt → approx px
            try:
                font = ImageFont.truetype("arial.ttf", fs_px)
            except Exception:
                try:
                    font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", fs_px)
                except Exception:
                    font = ImageFont.load_default()

            # Alignment
            align = para.alignment
            from pptx.enum.text import PP_ALIGN
            if align == PP_ALIGN.RIGHT:
                bbox = draw.textbbox((0, 0), line_text, font=font)
                text_w = bbox[2] - bbox[0]
                tx = x1 - text_w - 4
            elif align == PP_ALIGN.CENTER:
                bbox = draw.textbbox((0, 0), line_text, font=font)
                text_w = bbox[2] - bbox[0]
                tx = x0 + (x1 - x0 - text_w) // 2
            else:
                tx = x0 + 2

            draw.text((tx, cur_y), line_text, fill=line_color, font=font)
            bbox = draw.textbbox((tx, cur_y), line_text, font=font)
            cur_y = bbox[3] + int(line_size * SCALE * 0.15)

# Downscale to 1x for output
thumb = img.resize((720, 405), Image.LANCZOS)
thumb.save(OUT_PATH)
print(f"Thumbnail saved: {OUT_PATH}")
