# -*- coding: utf-8 -*-
"""Sets medium-length, simple speaker notes (~33 min, non-native friendly).
First sentence = key message, last sentence = transition."""
from pptx import Presentation

F = "research/presentation/China-Strategies-Presentation-FINAL.pptx"
prs = Presentation(F)

N = {
1: """Good afternoon, and welcome. Our paper asks one question: which way of building a car company in
China creates the most value for shareholders over time? We compare three real firms with three
models — SAIC uses joint ventures, BYD builds its own technology, and Geely buys companies abroad. Our
main point is simple: governance — who controls the firm — matters more than the technology it starts
with. We test this through four lenses: governance, financing, capability, and long-term viability.
Then we give our result, draw three lessons, and end with how we used AI. Each of us will take a part.
Let us start with why this question matters now.""",

2: """The joint-venture era is ending. Electric cars are now more than half of new sales in China, and
Chinese firms make about sixty percent of the world's EVs. For thirty years the deal was simple:
foreign firms gave technology, and a state firm gave market access. That worked when the engine and the
gearbox were the hard part. Now the hard part is the battery, the power electronics, and the software —
and that knowledge sits with the Chinese independents, not the foreign partners. So the old deal no
longer helps the joint-venture firms. You can see the stress at SAIC: its profit fell eighty-eight
percent in one year. And the EU set very different tariffs — high for SAIC, low for BYD — which tracks
how state-dependent each firm is. So the market and the regulator point the same way. Next, our exact
question.""",

3: """Our question is short but hard: which structure builds the most value over time? We study three
firms with three paths. SAIC is state-owned, controlled by the asset commission SASAC. BYD is
founder-led, run by Wang Chuanfu. Geely is also founder-led, under Li Shufu, but it grew by buying
firms abroad, like Volvo. To keep it fair, we look at all three through the same four lenses —
governance, financing, capability, and viability — in the same order. So we compare one structural
choice at a time, not three separate company stories. Our thesis is at the bottom: structure beats
starting technology. Now the theory behind these lenses.""",

4: """Six well-known ideas all point the same way. Shleifer and Vishny say that in a state firm,
control sends money to political goals, not to shareholders. Megginson and Netter show that private
firms beat state firms. Jensen and Meckling explain the good case: a strong founder-owner removes the
conflict between owners and managers. Holmes and Fan describe the joint-venture trade and its hidden
cost — connected firms lag about eighteen percent after listing. Luo and Tung explain why an
emerging-market firm buys abroad: to catch up quickly. And Teece explains why deep in-house knowledge
is so hard for rivals to copy. So the prediction is clear before we even look at the data: founder
control should adapt faster and create more value. Next, our map.""",

5: """This matrix is the map for the whole talk. The rows are the three firms, and the columns are the
four lenses. Read across a row for one firm's full profile; read down a column to compare the firms on
one topic. The key point is that the columns are linked. SAIC is state-led, capital-light, dependent on
transferred technology, and weakest in the EV shift — and those four cells belong together. BYD shows
the opposite pattern straight across, and Geely is in between. Over the next slides we go column by
column with the evidence behind each cell. We start with governance.""",

6: """Governance is really two questions: who holds control, and what they want. SAIC sits under the
state — SASAC owns about seventy-three percent, and the Party role is written into the charter, so its
goals go beyond profit, to things like jobs and policy. The cost is speed: it cannot leave a failing
joint venture without state approval. BYD is the opposite — the founder and his team hold over a third
of the votes, so Wang could end petrol-car production in one year, in 2022, on his own. Geely is also
founder-led and fast on deals, but it runs many brands, and some of them lose money. So control runs
on a scale, from political at SAIC to fully commercial at BYD. Next, the EU put a real price on this
difference.""",

7: """Here is proof from outside the firms. The EU studied state support and set tariffs: 35.3 percent
for SAIC, 17 for BYD, and 18.8 for Geely. That gap is not about car quality — it is about state
dependence. For SAIC the Commission even used “facts available,” because the group held back some
financing data. So a neutral regulator reached the same conclusion we reach from the inside: control
structure drives competitive position. And for a firm that wants to grow through exports, that tariff
is a direct, lasting cost of the governance model it chose. Now we follow the money.""",

8: """Each firm funds its growth differently, and the difference shows up in a downturn. SAIC was
capital-light: it shared cost and risk with its foreign partners. That felt safe — until SAIC-GM lost
more than half its sales, and SAIC was still bound to absorb the losses, with no control to fix the
business. Profit fell eighty-eight percent; the capital-light model had become a capital trap. BYD
funds itself from its own cash, and it spends more on research than it earns in profit, so it needs no
banks; the limit is its own cash flow. Geely borrows to buy capability fast — that is how it got
Volvo — but its debt is rising. So the real choice is resilience versus speed. Next, the numbers.""",

9: """Here are the numbers in one table. Look at return on equity: BYD and Geely are near twenty
percent, while SAIC is below one — that single line tells the story. Research spending: BYD spends more
on R&D than it earns in profit. Cash flow: BYD generates far more than the other two. One honest note
on SAIC: its profit looks like it recovered in 2025, but that is mostly because the big write-downs of
2024 did not repeat — the business stabilised, it did not really turn around. And the bottom row, the
EU tariff, marks the same governance gap once more. Next, how each firm builds its technology.""",

10: """How a firm learns shapes what it can build later. SAIC took technology through joint ventures.
It worked at first — quality rose about eight percent — but the easy money removed the pressure to
invent, so it had no EV platform of its own when the shift came. BYD built everything itself, from
battery cells to software, and that deep knowledge is very hard to copy. Geely bought Volvo and turned
it into real shared engineering, not one-way copying, and patent data confirm it. So the lesson is
simple: handed-over technology fades when the cycle turns, while built or co-developed knowledge lasts.
Next, the trade-off behind this.""",

11: """This slide shows the trade-off between speed, depth, and independence. SAIC was the fastest to
quality, but shallow and dependent, with no platform of its own. Geely is the rare strong case — fast
and deep — because it co-developed with Volvo. BYD was the slowest but the deepest, thirty years from
cells to software, and fully its own. The cost of the shallow path is in the bar at the bottom: SAIC's
joint-venture plants run at just 14 to 37 percent of capacity. When the technology shifts, there is no
platform to fall back on, and the factories sit half empty. Next, the EV transition itself.""",

12: """The EV shift tests every earlier choice. SAIC is the most exposed — its plants run far below
capacity — but to be fair, its own brand grew thirty-three percent, with no direct subsidies, and it
launched a brand with Huawei. So its future is a race between falling joint-venture income and a thin
own-brand margin, and its slow structure sets the pace. BYD is the strongest: no petrol cars since
2022, the world's number one for four years, and exports up one hundred and forty percent. Geely is in
between — strong growth, but it leans on partners and carries loss-making brands. Same shock, three
different results — exactly what the earlier choices predicted. Now our result.""",

13: """Here is our result. Of the three, BYD is the one most likely to create long-term value. The
evidence is consistent: the highest return on equity, four years as the top EV maker, and a
5.6-billion-dollar raise that sovereign wealth funds joined. But please notice the wording in the box.
It is not that building everything yourself is magic. BYD wins because it combines two things —
commercial founder control and a self-financing balance sheet that stays strong in a downturn. That is
governance and capital discipline working together. And we stay honest: BYD's own return fell to
thirteen percent in 2025 under the price war. Now three lessons that go beyond cars.""",

14: """We draw three lessons for any fast-moving market. One: how you build capability decides how long
it lasts. Buying or borrowing technology can hollow out a firm when the technology changes, while
building or co-developing keeps the strength. Two: governance sets the speed of change. Commercial
owners move fast; split control stalls — and when technology moves faster than your decisions,
governance becomes the limit. Three: capital structure decides who survives a downturn. Risk-sharing
can trap you, self-financing trades speed for safety, and debt does the opposite. One theme runs
through all three: structure decides. Now the conclusion.""",

15: """Here is the core point. Across all four lenses, structure — more than starting technology —
decided the outcome. The joint-venture model did not lose to a better battery; it lost to a better way
of allocating control and capital. BYD's founder control plus self-financing proved the most adaptive,
while SAIC's state model serves goals beyond profit, which shows up as a lasting discount. So our line,
in the box: when technology moves faster than firms can decide, governance is no longer a footnote to
strategy — it is strategy. Before questions, a short word on our limits and on how we used AI.""",

16: """A short, honest word on limits. We study three firms in one industry and one country, so this is
deep case work, not a statistical law. The EV shift is still young, so the ranking could change. We
rely on reported figures, and the reports mix currencies, which we reconciled carefully. We also left
NIO out — a fourth, capital-market model that was beyond our scope. Future work could add NIO, follow
these firms over a longer cycle, build a formal governance score, and study how the EU tariffs reshape
exports. Now, how we used AI.""",

17: """We did not use a simple chatbot. We used Claude Code as a layered system — a main model for most
work, a stronger one for the hard analysis, and five small agents, each with one job. But the heart of
it is the shared repository you see here. We had one folder for each case and a shared folder for
common data, like the EU regulation and the IEA figures. The rule was simple: pull before each
session, push after. Because the tool could read all three case wikis at once, we compared the firms
directly, with no manual file-sharing. Next, what this actually added.""",

18: """Where did AI help the most? First, it saved a lot of time — more than thirty papers and reports,
all ingested and searchable. Second, it kept us disciplined: every source was rated by quality, and we
built our arguments only on the strongest ones. The sharpest habit was to tell it to act like a strict
professor, and to quote the exact source for each fact. That same habit caught three real data errors
before submission — a wrong tariff, a profit figure inflated ten times, and an earnings figure on the
wrong share count — and we fixed all three against the audited reports. Next, the honest limits.""",

19: """Now the limits, honestly. It cannot download paywalled PDFs — it knows the papers' content, but
we still had to read and check the key claims ourselves. On long sessions we hit memory limits, and the
output was often too wordy and needed editing. The lesson is that this is not a chat box — it took time
to learn, and every output needed a careful read. The cost was small, about twenty-one euros a month
each. The line we held: AI was our tool and our reviewer, never a ghostwriter. The thinking stayed
ours. Finally, our sources.""",

20: """These are our main sources, in APA style. The theory comes from Shleifer and Vishny, Jensen and
Meckling, and Teece; the China evidence from Bai, Howell, and Holmes. The figures come from the audited
annual reports of all three firms, plus the EU regulation and the IEA outlook. The full list is in the
paper. Thank you very much for your attention — we are happy to take your questions.""",
}

for idx, txt in N.items():
    prs.slides[idx - 1].notes_slide.notes_text_frame.text = " ".join(txt.split())

prs.save(F)
print("notes updated; slides:", len(prs.slides))
