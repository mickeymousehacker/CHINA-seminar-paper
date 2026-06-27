"""
Build: "China Strategies in the Automotive Industry"
Seminar presentation — Goethe University Frankfurt, Session 6 (China strategies).

Design brief: McKinsey/BCG/Bain executive deck.
  - One slide = one core message in the headline (an assertion, not a topic).
  - China-EV palette: dark navy primary, electric blue, turquoise, red highlight.
  - Cambria Bold headlines, Calibri sub-heads/body. No Aptos, no Georgia.
  - No full-width colour bars, no accent stripe under titles.
  - Every slide carries a visual (stat callout, matrix, table, scorecard, quote).
  - Speaker notes on every slide: simple B2 English, short sentences.

Run:  uv run --with python-pptx python research/presentation/build_china_strategies.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x17, 0x2A)   # primary dark navy
BLUE   = RGBColor(0x25, 0x63, 0xEB)   # electric blue   → BYD
TEAL   = RGBColor(0x06, 0xB6, 0xD4)   # turquoise       → Geely
RED    = RGBColor(0xDC, 0x26, 0x26)   # China-red       → SAIC / highlight
BG     = RGBColor(0xF8, 0xFA, 0xFC)   # light background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x0F, 0x17, 0x2A)   # body text (= navy)
GREY   = RGBColor(0x64, 0x74, 0x8B)   # secondary text (slate-500)
LINE   = RGBColor(0xCB, 0xD5, 0xE1)   # hairlines / borders (slate-300)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)   # positive
AMBER  = RGBColor(0xD9, 0x77, 0x06)   # caution
NAVY2  = RGBColor(0x1E, 0x29, 0x3B)   # slightly lifted navy for cards on dark

# Company colours
SAIC   = RED
BYD    = BLUE
GEELY  = TEAL

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

W = Inches(13.333)
H = Inches(7.5)
MARGIN = Inches(0.6)

# ── Core helpers ──────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color, line_color=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    fill(sp, color)
    if line_color is not None:
        sp.line.color.rgb = line_color
        sp.line.width = line_w or Pt(0.75)
    return sp


def oval(slide, l, t, d, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    sp.shadow.inherit = False
    fill(sp, color)
    return sp


def text(slide, l, t, w, h, lines, size, color, *, font=BODY_FONT,
         bold=False, italic=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space_after=2, line_spacing=1.0):
    """lines: str (split on \\n) or list of (str, opts-dict)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = lines.split("\n")
    first = True
    for ln in lines:
        opts = {}
        if isinstance(ln, tuple):
            ln, opts = ln
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", align)
        p.space_after = Pt(opts.get("space_after", space_after))
        p.space_before = Pt(0)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        # a paragraph may hold mixed runs via opts["runs"] = [(txt, {...}), ...]
        runs = opts.get("runs")
        if runs is None:
            runs = [(ln, {})]
        for rtxt, ro in runs:
            r = p.add_run()
            r.text = rtxt
            r.font.name = ro.get("font", opts.get("font", font))
            r.font.size = Pt(ro.get("size", opts.get("size", size)))
            r.font.bold = ro.get("bold", opts.get("bold", bold))
            r.font.italic = ro.get("italic", opts.get("italic", italic))
            r.font.color.rgb = ro.get("color", opts.get("color", color))
    return tb


def bullets(slide, l, t, w, h, items, size, color=INK, *, gap=6,
            marker="–", marker_color=None, line_spacing=1.05):
    """items: list of str or (str, level)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        mk = ("   " if lvl else "") + (marker + "  " if marker else "")
        if mk:
            rm = p.add_run()
            rm.text = mk
            rm.font.name = BODY_FONT
            rm.font.size = Pt(size)
            rm.font.bold = False
            rm.font.color.rgb = marker_color or BLUE
        r = p.add_run()
        r.text = it
        r.font.name = BODY_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color if lvl == 0 else GREY
    return tb


def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt.strip()


def kicker_headline(slide, kicker, headline, kcolor=BLUE, hsize=27):
    """Standard light content-slide top: small kicker + Cambria headline."""
    rect(slide, 0, 0, W, H, BG)
    text(slide, MARGIN, Inches(0.42), Inches(12.1), Inches(0.3),
         kicker.upper(), 12.5, kcolor, font=BODY_FONT, bold=True)
    text(slide, MARGIN, Inches(0.72), Inches(12.1), Inches(1.0),
         headline, hsize, NAVY, font=HEAD_FONT, bold=True, line_spacing=1.0)


def footer(slide, page):
    text(slide, MARGIN, H - Inches(0.42), Inches(9.0), Inches(0.3),
         "China Strategies in the Automotive Industry  ·  Session 6",
         9, GREY, font=BODY_FONT)
    text(slide, W - Inches(1.6), H - Inches(0.42), Inches(1.0), Inches(0.3),
         str(page), 9, GREY, font=BODY_FONT, align=PP_ALIGN.RIGHT)


def chip(slide, l, t, w, label, color, h=Inches(0.4), size=12):
    rect(slide, l, t, w, h, color)
    text(slide, l, t + Inches(0.015), w, h, label, size, WHITE,
         font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


def card(slide, l, t, w, h, top_color, title, body_items, *,
         title_size=15, body_size=12, body_color=INK):
    """White card with a thin colour cap, a title, and bullet body."""
    rect(slide, l, t, w, h, CARD, line_color=LINE, line_w=Pt(0.75))
    rect(slide, l, t, w, Inches(0.09), top_color)
    text(slide, l + Inches(0.18), t + Inches(0.22), w - Inches(0.36), Inches(0.4),
         title, title_size, top_color, font=BODY_FONT, bold=True)
    bullets(slide, l + Inches(0.18), t + Inches(0.72),
            w - Inches(0.36), h - Inches(0.9), body_items, body_size,
            color=body_color, gap=5, marker="–", marker_color=top_color)


# ── Table helper ──────────────────────────────────────────────────────────────

def grid_table(slide, top, left, col_ws, row_h, headers, rows, *,
               head_bg=NAVY, head_fg=WHITE, font_size=12, head_size=12.5,
               first_col_fg=None, zebra=True, cell_colors=None):
    """col_ws: list of column widths. cell_colors: optional dict {(r,c):RGB}."""
    # header
    x = left
    for ci, hd in enumerate(headers):
        cw = col_ws[ci]
        rect(slide, x, top, cw, row_h, head_bg)
        text(slide, x + Inches(0.08), top, cw - Inches(0.16), row_h, hd,
             head_size, head_fg, font=BODY_FONT, bold=True,
             align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
        x += cw
    # rows
    for ri, row in enumerate(rows):
        y = top + row_h * (ri + 1)
        x = left
        base = WHITE if (ri % 2 == 0 or not zebra) else RGBColor(0xEE, 0xF2, 0xF7)
        for ci, cell in enumerate(row):
            cw = col_ws[ci]
            cc = base
            if cell_colors and (ri, ci) in cell_colors:
                cc = cell_colors[(ri, ci)]
            rect(slide, x, y, cw, row_h, cc, line_color=LINE, line_w=Pt(0.5))
            fg = INK
            if ci == 0 and first_col_fg:
                fg = first_col_fg
            text(slide, x + Inches(0.08), y, cw - Inches(0.16), row_h, cell,
                 font_size, fg, font=BODY_FONT,
                 bold=(ci == 0),
                 align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
            x += cw


def source_note(slide, txt, y=None):
    text(slide, MARGIN, y or (H - Inches(0.72)), Inches(12.1), Inches(0.3),
         txt, 9.5, GREY, font=BODY_FONT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def s01_title(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    # subtle vertical accent band on the left
    rect(s, 0, 0, Inches(0.16), H, BLUE)

    text(s, MARGIN, Inches(0.55), Inches(12.0), Inches(0.35),
         "M&A AS TRANSFORMER OF THE AUTOMOTIVE INDUSTRY", 13, TEAL,
         font=BODY_FONT, bold=True)

    text(s, MARGIN, Inches(2.0), Inches(12.0), Inches(1.7),
         "China Strategies in the\nAutomotive Industry", 46, WHITE,
         font=HEAD_FONT, bold=True, line_spacing=1.0)

    text(s, MARGIN, Inches(3.95), Inches(11.5), Inches(0.7),
         "Which governance and financing model creates long-term shareholder value?",
         18, RGBColor(0xCB, 0xD5, 0xE1), font=BODY_FONT, italic=True)

    # company tags
    tags = [("SAIC  ·  Joint Venture", SAIC),
            ("BYD  ·  Vertical Integration", BYD),
            ("Geely  ·  M&A-Driven", GEELY)]
    x = MARGIN
    for lbl, c in tags:
        chip(s, x, Inches(4.95), Inches(3.7), lbl, c, h=Inches(0.46), size=12.5)
        x += Inches(3.95)

    text(s, MARGIN, Inches(6.05), Inches(12.0), Inches(0.4),
         "Julia Ort  ·  Batughan Maus  ·  Chiara Stubner", 15, WHITE,
         font=BODY_FONT, bold=True)
    text(s, MARGIN, Inches(6.5), Inches(12.0), Inches(0.6),
         "Johann Wolfgang Goethe-Universität Frankfurt  ·  Chair of Banking and Finance\n"
         "Prof. Dr. Mark Wahrenburg  ·  Advisor: Edoardo Pilla  ·  29 June 2026",
         12, GREY, font=BODY_FONT, line_spacing=1.1)

    notes(s, """This presentation asks one question: which way of building a car company in China
creates the most value for shareholders over time. We study three real firms, each with a
different ownership and money model. SAIC uses joint ventures, BYD builds everything itself,
and Geely grows by buying companies. Our main finding is that governance — who controls the
firm — matters more than the technology a firm starts with. Let us begin with why this
question is urgent right now.""")


def s02_context(prs):
    s = blank(prs)
    kicker_headline(s, "Introduction · Market context",
                    "The 30-year joint-venture model is breaking down as EVs take over")
    stats = [
        ("~55%", "EV share of new\nvehicle sales in China (2025)", BLUE),
        ("~60%", "Chinese makers' share\nof global EV sales (2025)", TEAL),
        ("88%", "Drop in SAIC net\nprofit in 2024", RED),
        ("56.5%", "Fall in SAIC-GM\nsales in 2024", RED),
    ]
    cw = Inches(2.95)
    x = MARGIN
    for num, lbl, c in stats:
        rect(s, x, Inches(1.85), cw, Inches(2.0), CARD, line_color=LINE)
        rect(s, x, Inches(1.85), cw, Inches(0.09), c)
        text(s, x, Inches(2.2), cw, Inches(0.9), num, 40, c,
             font=HEAD_FONT, bold=True, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.12), Inches(3.12), cw - Inches(0.24), Inches(0.65),
             lbl, 12, INK, font=BODY_FONT, align=PP_ALIGN.CENTER, line_spacing=1.0)
        x += Inches(3.05)

    bullets(s, MARGIN, Inches(4.25), Inches(12.1), Inches(2.1), [
        "For three decades foreign technology was traded for market access, with a state-owned partner in the middle (Holmes et al., 2015) — that logic is now obsolete.",
        "Firms that leaned on foreign partners are losing share fast; independent Chinese makers have become global leaders (Whitfield & Wuttke, 2026).",
        "The EU now prices governance risk directly: countervailing duties of 35.3% on SAIC, 17.0% on BYD, 18.8% on Geely (European Commission, 2024).",
    ], 15, gap=9)
    source_note(s, "Sources: IEA Global EV Outlook 2026 · SAIC Annual Report 2024/2025 · European Commission (2024)")
    footer(s, 2)
    notes(s, """The joint-venture era is ending, and the numbers show it. Electric vehicles are now more
than half of all car sales in China, and Chinese firms make about 60 percent of the world's EVs.
SAIC, the classic joint-venture firm, saw profit fall 88 percent in one year. At the same time
the EU put the highest tariff on SAIC and the lowest on BYD. So the market and regulators both
point to the same thing: the old model is under pressure. Next, the exact question we ask and
how we study it.""")


def s03_question(prs):
    s = blank(prs)
    kicker_headline(s, "Introduction · Research design",
                    "One question, three governance archetypes, four analytical dimensions")
    # question box
    rect(s, MARGIN, Inches(1.75), Inches(12.1), Inches(0.95), NAVY)
    text(s, MARGIN + Inches(0.2), Inches(1.75), Inches(11.7), Inches(0.95),
         "“Which governance and financing structure is most likely to create long-term "
         "shareholder value in a technology-intensive, high-growth market?”",
         16, WHITE, font=HEAD_FONT, bold=True, italic=True,
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

    # three cases
    cases = [
        ("SAIC Motor", "Joint-venture model\nState-owned, SASAC ~73%", SAIC),
        ("BYD", "Vertical integration\nFounder-led, Wang ~17%", BYD),
        ("Geely", "M&A-driven\nFounder-led, Li ~41% attributable", GEELY),
    ]
    cw = Inches(3.93)
    x = MARGIN
    for title, desc, c in cases:
        rect(s, x, Inches(2.95), cw, Inches(1.35), CARD, line_color=LINE)
        rect(s, x, Inches(2.95), cw, Inches(0.09), c)
        text(s, x + Inches(0.18), Inches(3.13), cw - Inches(0.36), Inches(0.4),
             title, 16, c, font=BODY_FONT, bold=True)
        text(s, x + Inches(0.18), Inches(3.58), cw - Inches(0.36), Inches(0.7),
             desc, 12.5, INK, font=BODY_FONT, line_spacing=1.05)
        x += Inches(4.08)

    # four dimensions arrow strip
    text(s, MARGIN, Inches(4.55), Inches(12.1), Inches(0.3),
         "FOUR ANALYTICAL DIMENSIONS — APPLIED TO EACH CASE", 12, BLUE,
         font=BODY_FONT, bold=True)
    dims = ["1  Governance", "2  Financing", "3  Capability Building", "4  Long-term Viability"]
    dw = Inches(2.8)
    step = Inches(3.0)
    x = MARGIN
    for i, d in enumerate(dims):
        rect(s, x, Inches(4.95), dw, Inches(0.7), NAVY2)
        text(s, x + Inches(0.1), Inches(4.95), dw - Inches(0.2), Inches(0.7),
             d, 13.5, WHITE, font=BODY_FONT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            text(s, x + dw, Inches(4.95), step - dw, Inches(0.7),
                 "›", 18, GREY, font=BODY_FONT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += step

    text(s, MARGIN, Inches(5.95), Inches(12.1), Inches(0.6),
         "Central thesis: governance structure, more than initial technological endowment, "
         "determines long-term value creation.",
         14, NAVY, font=BODY_FONT, bold=True, line_spacing=1.05)
    footer(s, 3)
    notes(s, """Our question is simple to say but hard to answer: which structure builds the most
shareholder value over time. We compare three firms that each picked a different path — joint
ventures, building in-house, and buying companies. We look at each one through four lenses:
governance, financing, capability building, and long-term survival. The same four lenses are
used for all three firms, so the comparison is fair. Our thesis is that structure beats starting
technology. Now the theory behind these lenses.""")


def s04_theory(prs):
    s = blank(prs)
    kicker_headline(s, "Theory & literature",
                    "Six established lenses frame governance as the driver of value")
    th = [
        ("Control rights drive allocation", "Shleifer & Vishny (1994)",
         "State control steers firms toward political goals, away from value maximisation.", RED),
        ("State vs. private ownership", "Megginson & Netter (2001)",
         "Privatised firms outperform; mixed-ownership firms do not catch up.", RED),
        ("Agency & founder alignment", "Jensen & Meckling (1976)",
         "Concentrated founder ownership removes the principal–agent conflict.", BLUE),
        ("Quid pro quo transfer", "Holmes et al. (2015); Fan et al. (2007)",
         "Tech-for-access transfer; politically connected firms lag ~18% post-IPO.", RED),
        ("Springboard M&A", "Luo & Tung (2007)",
         "Emerging-market firms acquire abroad to leap over latecomer disadvantages.", TEAL),
        ("Dynamic capabilities", "Teece et al. (1997)",
         "Path-dependent tacit knowledge is the source of non-imitable advantage.", BLUE),
    ]
    cw = Inches(3.93)
    ch = Inches(2.15)
    gx = Inches(0.165)
    gy = Inches(0.22)
    for i, (t, cite, body, c) in enumerate(th):
        col = i % 3
        rowi = i // 3
        x = MARGIN + col * (cw + gx)
        y = Inches(1.85) + rowi * (ch + gy)
        rect(s, x, y, cw, ch, CARD, line_color=LINE)
        rect(s, x, y, Inches(0.09), ch, c)
        text(s, x + Inches(0.22), y + Inches(0.18), cw - Inches(0.36), Inches(0.55),
             t, 14.5, NAVY, font=BODY_FONT, bold=True, line_spacing=0.95)
        text(s, x + Inches(0.22), y + Inches(0.82), cw - Inches(0.36), Inches(0.3),
             cite, 11.5, c, font=BODY_FONT, bold=True, italic=True)
        text(s, x + Inches(0.22), y + Inches(1.18), cw - Inches(0.36), Inches(0.85),
             body, 12, INK, font=BODY_FONT, line_spacing=1.0)
    footer(s, 4)
    notes(s, """Our argument stands on six well-known ideas. Shleifer and Vishny and also Megginson
and Netter explain why state firms tend to lag private ones. Jensen and Meckling show why a
strong founder owner can act fast and align with shareholders. Holmes and Fan describe the
quid-pro-quo trade and its costs for connected firms. Luo and Tung explain why a firm buys
abroad to catch up, and Teece explains why deep in-house knowledge is hard to copy. Together
they predict that governance shapes who adapts. Next we turn these lenses into one framework.""")


def s05_framework(prs):
    s = blank(prs)
    kicker_headline(s, "Theory & literature · Framework",
                    "Three models tested against four dimensions — a 3 × 4 matrix")
    headers = ["", "Governance", "Financing", "Capability", "Viability"]
    rows = [
        ["", "State / SASAC + Party", "JV dividends\n(capital-light)", "Mandatory transfer", "Most exposed"],
        ["", "Founder-led\n(Wang ~17%)", "Self-financing\n(operating cash)", "Organic build", "Strongest"],
        ["", "Founder-led\n(Li ~41%)", "Leveraged M&A", "Acquisition + co-dev", "Most asymmetric"],
    ]
    col_ws = [Inches(1.7), Inches(2.85), Inches(2.55), Inches(2.55), Inches(2.45)]
    cell_colors = {
        (0, 0): RED, (1, 0): BLUE, (2, 0): TEAL,
    }
    # render with coloured first-column labels
    top = Inches(1.95)
    row_h = Inches(0.92)
    grid_table(s, top, MARGIN, col_ws, row_h, headers, rows,
               font_size=12.5, head_size=13, cell_colors=cell_colors)
    # white text on coloured label cells (overlay)
    labels = [("SAIC", RED), ("BYD", BLUE), ("Geely", TEAL)]
    for ri, (lbl, c) in enumerate(labels):
        y = top + row_h * (ri + 1)
        text(s, MARGIN, y, col_ws[0], row_h, lbl, 14, WHITE,
             font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    text(s, MARGIN, Inches(5.95), Inches(12.1), Inches(0.6),
         "Reading the matrix: the same firm that is least adaptable in governance is also the "
         "most exposed in viability — the columns are linked, not independent.",
         13.5, NAVY, font=BODY_FONT, bold=True, line_spacing=1.05)
    source_note(s, "Theoretical anchors: Shleifer & Vishny (1994) · Megginson & Netter (2001) · "
                   "Teece et al. (1997) · Luo & Tung (2007) · Fan et al. (2007)")
    footer(s, 5)
    notes(s, """This matrix is the map for the whole talk. The rows are the three firms. The columns
are the four lenses. Read across a row and you see one firm's full profile; read down a column
and you compare the firms on one topic. The key point is that the columns are linked. The firm
that is slow to decide is also the one most at risk in the EV shift. The next eight slides walk
through the four columns, two slides each. We start with governance.""")


def s06_governance(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 1 · Governance",
                    "Who controls the firm — and what they optimise — diverges sharply")
    card(s, MARGIN, Inches(1.85), Inches(3.93), Inches(4.2), SAIC, "SAIC — Dual state principal", [
        "SASAC owns ~73%; CCP mandate is written into the charter",
        "Optimises beyond shareholder value: jobs, policy, prestige",
        "Politically connected firms lag ~18% post-IPO (Fan et al., 2007)",
        "Cannot exit unprofitable JVs without state sign-off",
        "Slow capital reallocation → strategic rigidity",
    ], body_size=12)
    card(s, MARGIN + Inches(4.08), Inches(1.85), Inches(3.93), Inches(4.2), BYD, "BYD — Founder control", [
        "Wang Chuanfu ~17%; founding team >1/3 of votes",
        "No JV dependency; one clear commercial objective",
        "Founder control aids breakthrough innovation (Zhang et al., 2025)",
        "ICE exit (2022) decided unilaterally, no state approval",
        "Fastest decision cycle of the three",
    ], body_size=12)
    card(s, MARGIN + Inches(8.16), Inches(1.85), Inches(3.93), Inches(4.2), GEELY, "Geely — Founder apex, multi-brand", [
        "Li Shufu ~41% attributable via ZGH (64.4% combined)",
        "No state committee or JV-partner consent at the top",
        "Agile in M&A, but coordination cost across brands",
        "Volvo / Polestar / Lotus add complexity",
        "Loss-making subsidiaries = contingent liability",
    ], body_size=12)
    source_note(s, "Theory: Shleifer & Vishny (1994) · Fan et al. (2007) · Jensen & Meckling (1976) · Megginson & Netter (2001)")
    footer(s, 6)
    notes(s, """Governance means two things: who holds control, and what they try to maximise. SAIC
is controlled by the state, so it must serve goals beyond profit, and it cannot move fast. BYD
is run by its founder, who owns a big block and can decide alone, as he did when he stopped
making petrol cars in 2022. Geely is also founder-led and quick at deals, but managing many
brands adds friction and risk. So control sits on a scale from political to commercial. The
next slide shows that the EU put a price on exactly this difference.""")


def s07_governance_synth(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 1 · Governance synthesis",
                    "The EU tariff gap prices governance risk — a third-party validation")
    tariffs = [("35.3%", "SAIC", "State-owned;\nhighest subsidy finding", RED),
               ("17.0%", "BYD", "Founder-led;\nlowest state finding", BLUE),
               ("18.8%", "Geely", "Private;\nmoderate state finding", TEAL)]
    cw = Inches(3.6)
    x = MARGIN + Inches(0.35)
    for rate, name, note_, c in tariffs:
        rect(s, x, Inches(1.85), cw, Inches(2.2), CARD, line_color=LINE)
        rect(s, x, Inches(1.85), cw, Inches(0.09), c)
        text(s, x, Inches(2.1), cw, Inches(0.95), rate, 44, c,
             font=HEAD_FONT, bold=True, align=PP_ALIGN.CENTER)
        text(s, x, Inches(3.05), cw, Inches(0.4), name, 16, NAVY,
             font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.12), Inches(3.45), cw - Inches(0.24), Inches(0.55),
             note_, 11.5, GREY, font=BODY_FONT, italic=True,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
        x += Inches(3.85)

    bullets(s, MARGIN, Inches(4.4), Inches(12.1), Inches(1.9), [
        "An independent regulator — the European Commission — priced SAIC's state dependence at an 18.3 pp premium over BYD.",
        "Principal-agent hierarchy: JV boards (bilateral consent) > SOE oversight > founder control — adaptability rises down the chain.",
        "Core message: governance is a first-order determinant of strategic adaptability, not an administrative footnote.",
    ], 15, gap=9)
    source_note(s, "European Commission (2024), Implementing Regulation (EU) 2024/2754 · Megginson & Netter (2001)")
    footer(s, 7)
    notes(s, """Here is outside proof that governance matters. The EU studied state support and set
duties: 35.3 percent for SAIC, 17 for BYD, 18.8 for Geely. That gap of about 18 points is not
about car quality — it is about ownership. The more a firm depends on the state, the higher the
tariff. This lines up with our theory: control split across political and JV partners adapts
slower than founder control. Governance is now a direct cost, not paperwork. Next we follow the
money: how each firm pays for growth.""")


def s08_financing(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 2 · Financing",
                    "Three financing models, three very different risk profiles")
    headers = ["Metric (FY24 / FY25)", "SAIC", "BYD", "Geely"]
    rows = [
        ["Gross margin", "9.4% / 10.1%", "19.4% / 17.7%", "15.9% / 16.6%"],
        ["Net profit (RMB bn)", "1.7 / 10.1", "40.3 / 32.6", "16.8 / 16.9"],
        ["Return on equity", "0.58% / 3.43%", "21.7% / 13.2%", "19.2% / 18.2%"],
        ["R&D expense (RMB bn)", "35.2 / 33.6", "54.2 / 63.4", "10.4 / 17.6"],
        ["Operating cash flow (RMB bn)", "69.3 / 34.3", "133.5 / 59.1", "26.5 / 47.3"],
        ["Net gearing", "n.a. / n.a.", "–36% / +25%", "8.8% / 19.8%"],
        ["EU anti-subsidy tariff", "35.3%", "17.0%", "18.8%"],
    ]
    col_ws = [Inches(4.0), Inches(2.7), Inches(2.7), Inches(2.7)]
    row_h = Inches(0.55)
    grid_table(s, Inches(1.75), MARGIN, col_ws, row_h, headers, rows,
               font_size=12.5, head_size=13)
    # colour the company header cells
    for i, c in enumerate([SAIC, BYD, GEELY]):
        x = MARGIN + col_ws[0] + sum(col_ws[1:1+i], Inches(0))
        rect(s, x, Inches(1.75), col_ws[1+i], row_h, c)
        text(s, x, Inches(1.75), col_ws[1+i], row_h, ["SAIC", "BYD", "Geely"][i],
             13, WHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    text(s, MARGIN, Inches(6.3), Inches(12.1), Inches(0.45),
         "BYD self-funds (R&D > net profit); SAIC leans on JV dividends that collapsed in 2024; "
         "Geely funds capability with debt (gearing rising).",
         13, NAVY, font=BODY_FONT, bold=True, line_spacing=1.0)
    source_note(s, "Table 1, Section 3.4 · Annual Reports SAIC/BYD/Geely 2024 & 2025 · EU Reg. 2024/2754")
    footer(s, 8)
    notes(s, """This is the money picture, taken straight from the paper's Table 1. Look at three
rows. Return on equity: BYD and Geely are near 20 percent, SAIC is below one. Research spending:
BYD spends more on R&D than it earns in profit, which is rare. Cash flow: BYD generates far more
cash from operations than the others. SAIC's model was capital-light but turned into a capital
trap when JV income fell. Geely buys speed with debt, so its gearing is climbing. Next, what
this means for risk.""")


def s09_financing_synth(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 2 · Risk allocation",
                    "Capital structure decides who survives the downturn")
    cols = [
        ("Capital-light JV", SAIC,
         ["SAIC externalised investment risk via 50/50 JVs",
          "Worked while markets grew — became a capital trap",
          "SAIC-GM sales –56.5%; plants at 14–37% utilisation",
          "Bound to absorb losses without control to restructure"]),
        ("Self-financing", BYD,
         ["Funded from operating cash, not external lenders",
          "R&D > net profit; net-cash position for resilience",
          "USD 5.6 bn raise with sovereign wealth funds",
          "Resilient — but growth capped by internal cash"]),
        ("Leveraged M&A", GEELY,
         ["Debt-financed acquisitions buy capability fast",
          "Gearing rose to 19.8% (FY25), peak 22.2% in H1",
          "Operating cash flow RMB 47.3 bn gives headroom",
          "Speed — at the cost of balance-sheet exposure"]),
    ]
    cw = Inches(3.93)
    x = MARGIN
    for title, c, items in cols:
        card(s, x, Inches(1.85), cw, Inches(3.65), c, title, items, body_size=12)
        x += Inches(4.08)

    rect(s, MARGIN, Inches(5.75), Inches(12.1), Inches(0.85), NAVY)
    text(s, MARGIN + Inches(0.2), Inches(5.75), Inches(11.7), Inches(0.85),
         "Capital-light captures upside in growth but traps the firm in a downturn; self-financing "
         "trades growth speed for resilience; leverage trades balance-sheet safety for strategic speed.",
         13.5, WHITE, font=BODY_FONT, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    footer(s, 9)
    notes(s, """Each money model handles a downturn differently. SAIC shared risk with partners, which
felt safe, but when sales fell it still had to absorb losses it could not control — a capital
trap. BYD pays for itself, so it does not depend on banks; the cost is slower growth. Geely
borrows to move fast, which works while cash flow is strong but adds risk if rates or sales
turn. So the choice is really about resilience versus speed. Mixed-ownership firms tend to lag
private peers. Now the third lens: how each firm builds know-how.""")


def s10_capability(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 3 · Capability building",
                    "How a firm learns determines what it can do next")
    cols = [
        ("SAIC — Mandatory transfer", SAIC,
         ["Tech-for-access via state-mandated JVs",
          "+8.3% quality from quid pro quo (Bai et al., 2025)",
          "But JV rents suppressed own R&D — “like opium” (Howell, 2018)",
          "No proprietary EV platform when the cycle turned",
          "Capability paradox: scale without sovereignty"]),
        ("BYD — Organic build", BYD,
         ["Battery → power electronics → software, in-house",
          "Built ahead of state policy (Whitfield & Wuttke, 2026)",
          "Tacit, path-dependent knowledge — non-imitable (Teece et al., 1997)",
          "Slowest to build, but deepest and most durable",
          "Full vertical stack = structural moat"]),
        ("Geely — Acquisition + co-dev", GEELY,
         ["Volvo (2010): brand and engineering depth at once",
          "CMA platform = genuine cross-brand co-development",
          "Patent convergence confirms co-creation (Konda et al., 2022)",
          "Fastest to global tech parity of the three",
          "Risk: depends on arm's-length partner relations"]),
    ]
    cw = Inches(3.93)
    x = MARGIN
    for title, c, items in cols:
        card(s, x, Inches(1.85), cw, Inches(4.2), c, title, items, body_size=11.5)
        x += Inches(4.08)
    source_note(s, "Bai et al. (2025) · Howell (2018) · Teece et al. (1997) · Whitfield & Wuttke (2026) · Konda et al. (2022)")
    footer(s, 10)
    notes(s, """How a firm learns shapes what it can build later. SAIC got quality from its foreign
partners — about 8 percent better cars — but the easy JV money killed the urge to invent its
own technology, so it had no EV platform ready. BYD learned slowly and by itself, from batteries
up to software; that knowledge is deep and very hard to copy. Geely bought Volvo and turned the
deal into real shared engineering, like the CMA platform. So we have three speeds and three
depths of learning. The next slide shows the trade-off and what it cost.""")


def s11_capability_synth(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 3 · Speed vs. depth",
                    "Capability building trades speed against depth against independence")
    # 2x2 matrix: x = speed (slow→fast), y = depth/independence (low→high)
    ox, oy = MARGIN, Inches(1.95)
    mw, mh = Inches(7.4), Inches(4.1)
    rect(s, ox, oy, mw, mh, WHITE, line_color=LINE)
    # axes labels
    text(s, ox, oy + mh + Inches(0.05), mw, Inches(0.3),
         "SPEED OF ACQUISITION  →", 11, GREY, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    tb = text(s, ox - Inches(1.7), oy + mh/2 - Inches(0.16), Inches(2.9), Inches(0.32),
              "DEPTH / INDEPENDENCE  →", 10.5, GREY, font=BODY_FONT, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb.rotation = 270
    # gridlines
    rect(s, ox + mw/2, oy, Pt(1), mh, LINE)
    rect(s, ox, oy + mh/2, mw, Pt(1), LINE)
    # plot points
    def plot(cx, cy, label, c):
        d = Inches(0.34)
        oval(s, cx - d/2, cy - d/2, d, c)
        text(s, cx + Inches(0.22), cy - Inches(0.18), Inches(2.6), Inches(0.5),
             label, 12.5, NAVY, font=BODY_FONT, bold=True)
    plot(ox + Inches(5.3), oy + Inches(3.25), "SAIC", SAIC)  # fast entry, shallow & dependent
    plot(ox + Inches(1.4), oy + Inches(0.95), "BYD", BYD)    # slow, deep & independent
    plot(ox + Inches(5.1), oy + Inches(1.55), "Geely", GEELY)  # fast, medium-deep
    # quadrant hints
    text(s, ox + Inches(0.15), oy + Inches(0.12), Inches(3.4), Inches(0.3),
         "slow · deep & independent", 10, GREY, font=BODY_FONT, italic=True)
    text(s, ox + mw - Inches(3.5), oy + Inches(0.12), Inches(3.4), Inches(0.3),
         "fast · deep (rare)", 10, GREY, font=BODY_FONT, italic=True, align=PP_ALIGN.RIGHT)
    text(s, ox + Inches(0.15), oy + mh - Inches(0.32), Inches(3.4), Inches(0.3),
         "slow · shallow", 10, GREY, font=BODY_FONT, italic=True)
    text(s, ox + mw - Inches(3.5), oy + mh - Inches(0.32), Inches(3.4), Inches(0.3),
         "fast · shallow & dependent", 10, GREY, font=BODY_FONT, italic=True, align=PP_ALIGN.RIGHT)

    # right side takeaways
    rx = ox + mw + Inches(0.4)
    rw = Inches(4.1)
    text(s, rx, oy, rw, Inches(0.35), "THE PROOF IN UTILISATION", 12, BLUE,
         font=BODY_FONT, bold=True)
    for i, (num, lbl, c) in enumerate([
        ("37%", "SAIC-GM capacity utilisation (2025)", RED),
        ("14%", "VW MEB plant utilisation (2025)", RED),
    ]):
        y = oy + Inches(0.45) + i * Inches(1.0)
        rect(s, rx, y, rw, Inches(0.9), CARD, line_color=LINE)
        text(s, rx + Inches(0.15), y + Inches(0.1), Inches(1.4), Inches(0.7),
             num, 30, c, font=HEAD_FONT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        text(s, rx + Inches(1.6), y, rw - Inches(1.7), Inches(0.9),
             lbl, 12, INK, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text(s, rx, oy + Inches(2.55), rw, Inches(1.5),
         "Mandatory transfer (SAIC) built scale fast but shallow — dependency hardened into a "
         "capability deficit. Organic depth (BYD) and co-development (Geely) preserved autonomy.",
         12.5, NAVY, font=BODY_FONT, line_spacing=1.1)
    footer(s, 11)
    notes(s, """This chart sums up learning as a trade-off. The right means fast, the top means deep
and independent. SAIC sits low: it learned fast but stayed shallow and dependent. BYD sits high
but to the left: slow yet deep. Geely is in the middle-right: fairly fast and fairly deep. The
cost of the shallow path shows in the numbers on the right — SAIC's joint-venture plants run at
just 14 to 37 percent. That is dependency turning into a real deficit. Now the fourth lens: the
EV transition as a test of all these choices.""")


def s12_viability(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 4 · Long-term viability",
                    "The EV transition is a structural audit of past governance choices")
    cols = [
        ("SAIC — Most exposed", SAIC,
         ["NEV sales 1.643 m units (+33%) — real progress",
          "But a race vs. falling JV income and thin margins",
          "Highest EU tariff (35.3%) blocks the export escape",
          "Recovery rests on absent impairments, not operations",
          "Turnaround needs state support or JV renegotiation"]),
        ("BYD — Strongest", BYD,
         ["No ICE since 2022; #1 NEV maker 4 years running",
          ">1 m exports (+140%), present in 119 countries",
          "USD 5.6 bn raise with sovereign wealth funds",
          "Self-funded EV scale-up; lowest tariff (17%)",
          "Margin pressure: gross margin 19.4% → 17.7%"]),
        ("Geely — Most asymmetric", GEELY,
         ["NEV share 51.5% in H1 2025 (+126%)",
          "3.02 m deliveries FY2025 (+39%); H1 profit +102%",
          "Operating cash flow RMB 47.3 bn gives headroom",
          "But platforms depend on Volvo / Smart JV partners",
          "Polestar & Lotus losses = contingent liability"]),
    ]
    cw = Inches(3.93)
    x = MARGIN
    for title, c, items in cols:
        card(s, x, Inches(1.85), cw, Inches(4.2), c, title, items, body_size=11.5)
        x += Inches(4.08)
    source_note(s, "Annual Reports 2025 (FY2024) & 2026 (FY2025) · Geely interim results H1 2025 · European Commission (2024)")
    footer(s, 12)
    notes(s, """The EV shift acts like an exam for past decisions. SAIC is most exposed: it is
growing EV sales, but it is racing against falling joint-venture income, and the high EU tariff
closes the export door. BYD is strongest: no petrol cars since 2022, world number one in NEVs,
and able to fund its own growth. Geely is the mixed case: very strong sales growth, but it leans
on partners for platforms and carries loss-making brands. So the same shock produces three very
different results. The next slide scores them side by side.""")


def s13_scorecard(prs):
    s = blank(prs)
    kicker_headline(s, "Dimension 4 · Structural audit",
                    "BYD enters the EV era strongest; SAIC the most exposed")
    crit = ["EV readiness", "Margin resilience", "Balance-sheet", "Adaptability", "Overall"]
    # rating: 2=green,1=amber,0=red
    data = {
        "SAIC": [0, 0, 1, 0, 0],
        "BYD":  [2, 1, 2, 2, 2],
        "Geely":[2, 1, 1, 1, 1],
    }
    colormap = {2: GREEN, 1: AMBER, 0: RED}
    labelmap = {2: "Strong", 1: "Mixed", 0: "Weak"}
    # header row
    top = Inches(2.05)
    lab_w = Inches(1.9)
    cw = Inches(2.04)
    x = MARGIN + lab_w
    rect(s, MARGIN, top, lab_w, Inches(0.6), NAVY)
    for c in crit:
        rect(s, x, top, cw, Inches(0.6), NAVY)
        text(s, x + Inches(0.05), top, cw - Inches(0.1), Inches(0.6), c, 11.5, WHITE,
             font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.9)
        x += cw
    # rows
    comp_colors = {"SAIC": SAIC, "BYD": BYD, "Geely": GEELY}
    for ri, comp in enumerate(["SAIC", "BYD", "Geely"]):
        y = top + Inches(0.6) + ri * Inches(1.0)
        rect(s, MARGIN, y, lab_w, Inches(1.0), comp_colors[comp])
        text(s, MARGIN, y, lab_w, Inches(1.0), comp, 16, WHITE,
             font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = MARGIN + lab_w
        for ci in range(len(crit)):
            rect(s, x, y, cw, Inches(1.0), WHITE, line_color=LINE)
            val = data[comp][ci]
            d = Inches(0.32)
            oval(s, x + cw/2 - d/2, y + Inches(0.18), d, colormap[val])
            text(s, x, y + Inches(0.55), cw, Inches(0.35), labelmap[val], 10.5,
                 colormap[val], font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
            x += cw

    text(s, MARGIN, Inches(6.1), Inches(12.1), Inches(0.7),
         "Architecture decides EV-transition speed: founder control + self-financing (BYD) "
         "converts the technology shock into leadership; state control + JV dependence (SAIC) "
         "converts it into a structural discount.",
         13.5, NAVY, font=BODY_FONT, bold=True, line_spacing=1.05)
    # legend
    lx = W - Inches(4.6)
    for i, (v, lbl) in enumerate([(2, "Strong"), (1, "Mixed"), (0, "Weak")]):
        oval(s, lx + i * Inches(1.5), H - Inches(0.46), Inches(0.18), colormap[v])
        text(s, lx + i * Inches(1.5) + Inches(0.24), H - Inches(0.5), Inches(1.2), Inches(0.3),
             lbl, 10, GREY, font=BODY_FONT)
    footer(s, 13)
    notes(s, """This scorecard puts the audit on one page. Green means strong, amber mixed, red weak.
BYD is green almost everywhere — ready for EVs, strong balance sheet, fast to adapt. SAIC is
mostly red, especially on adaptability, because the state structure slows it down. Geely sits in
between, strong on EV products but only mixed on its balance sheet. The pattern is clear: the
structure a firm chose years ago decides how fast it can move now. From here we draw the general
lessons.""")


def s14_lessons12(prs):
    s = blank(prs)
    kicker_headline(s, "Generalizable lessons",
                    "Lessons 1 & 2: capability mode and governance set the limits of adaptation")
    # Lesson 1
    rect(s, MARGIN, Inches(1.85), Inches(5.95), Inches(4.2), CARD, line_color=LINE)
    rect(s, MARGIN, Inches(1.85), Inches(5.95), Inches(0.09), BLUE)
    text(s, MARGIN + Inches(0.2), Inches(2.05), Inches(5.55), Inches(0.3),
         "LESSON 1", 12, BLUE, font=BODY_FONT, bold=True)
    text(s, MARGIN + Inches(0.2), Inches(2.35), Inches(5.55), Inches(0.8),
         "Capability mode determines durability across technology cycles", 16, NAVY,
         font=BODY_FONT, bold=True, line_spacing=1.0)
    bullets(s, MARGIN + Inches(0.2), Inches(3.35), Inches(5.55), Inches(2.5), [
        "Mandatory transfer buys short-run quality but suppresses the autonomous innovation needed for the next cycle.",
        "Organic build and voluntary M&A preserve that innovation — trading speed against depth against independence.",
        "SAIC mastered assembly but held no battery, power-electronics or software platform when EVs arrived.",
    ], 13, gap=8)

    # Lesson 2
    lx = MARGIN + Inches(6.15)
    rect(s, lx, Inches(1.85), Inches(5.95), Inches(4.2), CARD, line_color=LINE)
    rect(s, lx, Inches(1.85), Inches(5.95), Inches(0.09), RED)
    text(s, lx + Inches(0.2), Inches(2.05), Inches(5.55), Inches(0.3),
         "LESSON 2", 12, RED, font=BODY_FONT, bold=True)
    text(s, lx + Inches(0.2), Inches(2.35), Inches(5.55), Inches(0.8),
         "Governance is a first-order determinant of adaptability", 16, NAVY,
         font=BODY_FONT, bold=True, line_spacing=1.0)
    bullets(s, lx + Inches(0.2), Inches(3.35), Inches(5.55), Inches(1.7), [
        "Speed of capital reallocation depends on how control rights are distributed.",
        "Control with commercial principals adapts faster than control split across political or JV partners.",
    ], 13, gap=8)
    # mini contrast
    rect(s, lx + Inches(0.2), Inches(5.05), Inches(2.65), Inches(0.85), RGBColor(0xFE, 0xE2, 0xE2))
    text(s, lx + Inches(0.3), Inches(5.1), Inches(2.45), Inches(0.75),
         "SAIC, 2022: ICE exit\nstructurally impossible", 12, RED, font=BODY_FONT,
         bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    rect(s, lx + Inches(3.0), Inches(5.05), Inches(2.75), Inches(0.85), RGBColor(0xDB, 0xEA, 0xFE))
    text(s, lx + Inches(3.1), Inches(5.1), Inches(2.55), Inches(0.75),
         "BYD, 2022: full ICE exit\ndecided in one year", 12, BLUE, font=BODY_FONT,
         bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    footer(s, 14)
    notes(s, """Two general lessons come out of the cases. First, how a firm builds skills decides how
long it stays competitive. Forced transfer gives quick quality but leaves a firm hollow when
the technology changes, as SAIC found with EVs. Second, governance sets the speed of change.
The clearest proof is 2022: BYD's founder could quit petrol cars in a single year, while SAIC's
state structure made the same move almost impossible. So control rights, not just resources,
decide who adapts. The last lesson and our conclusion come next.""")


def s15_lesson3_concl(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, BLUE)
    text(s, MARGIN, Inches(0.45), Inches(12.1), Inches(0.3),
         "GENERALIZABLE LESSONS · CONCLUSION", 12.5, TEAL, font=BODY_FONT, bold=True)
    text(s, MARGIN, Inches(0.78), Inches(12.1), Inches(0.6),
         "Lesson 3 & conclusion: governance is strategy, not its footnote", 26, WHITE,
         font=HEAD_FONT, bold=True)

    # Lesson 3
    text(s, MARGIN, Inches(1.75), Inches(7.3), Inches(0.35),
         "LESSON 3", 12, TEAL, font=BODY_FONT, bold=True)
    text(s, MARGIN, Inches(2.05), Inches(7.3), Inches(0.5),
         "Capital structure determines resilience when markets contract", 16, WHITE,
         font=BODY_FONT, bold=True, line_spacing=1.0)
    bullets(s, MARGIN, Inches(2.9), Inches(7.3), Inches(2.0), [
        "Risk-sharing captures upside in growth but exposes the firm in downturns it cannot control.",
        "Mixed-ownership firms structurally underperform private peers (Megginson & Netter, 2001).",
        "Founder-led vertical integration is the strongest model for long-term shareholder value.",
    ], 13.5, color=RGBColor(0xE2, 0xE8, 0xF0), gap=8, marker_color=TEAL)

    # BYD answer callout
    rect(s, MARGIN, Inches(5.05), Inches(7.3), Inches(1.05), NAVY2)
    text(s, MARGIN + Inches(0.2), Inches(5.05), Inches(7.0), Inches(1.05),
         [("Answer: ", {"bold": True, "color": WHITE}),
          ("BYD — ROE 21.7% (FY2024), 4 years #1 in NEVs, self-financed and adaptable.",
           {"color": RGBColor(0xE2, 0xE8, 0xF0)})],
         13.5, WHITE, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

    # quote panel
    qx = MARGIN + Inches(7.7)
    rect(s, qx, Inches(1.75), Inches(4.4), Inches(4.35), NAVY2)
    rect(s, qx, Inches(1.75), Inches(4.4), Inches(0.09), TEAL)
    text(s, qx + Inches(0.3), Inches(2.1), Inches(3.8), Inches(0.8),
         "“", 60, TEAL, font=HEAD_FONT, bold=True)
    text(s, qx + Inches(0.3), Inches(2.75), Inches(3.8), Inches(3.0),
         "When technology moves faster than firms can decide, governance is no longer a "
         "footnote to strategy — it is strategy.",
         18, WHITE, font=HEAD_FONT, bold=True, italic=True, line_spacing=1.1)
    footer(s, 15)
    notes(s, """The third lesson is about money under stress: the structure that captures gains in
good times can trap a firm in bad times, and self-financing gives the best resilience. Putting
all three lessons together gives our answer. BYD creates the most shareholder value: highest
return on equity, four years as the top EV maker, and able to fund itself. The reason is not
that vertical integration is magic. It is that BYD combines commercial founder control with a
resilient balance sheet. In short: governance is strategy. Next, the limits of our study.""")


def s16_limitations(prs):
    s = blank(prs)
    kicker_headline(s, "Limitations & future research",
                    "Scope conditions and where the analysis could be extended")
    rect(s, MARGIN, Inches(1.9), Inches(5.95), Inches(3.9), CARD, line_color=LINE)
    rect(s, MARGIN, Inches(1.9), Inches(5.95), Inches(0.09), AMBER)
    text(s, MARGIN + Inches(0.2), Inches(2.1), Inches(5.55), Inches(0.4),
         "Limitations", 16, AMBER, font=BODY_FONT, bold=True)
    bullets(s, MARGIN + Inches(0.2), Inches(2.6), Inches(5.55), Inches(3.0), [
        "Three firms in one industry and one country — limited statistical generalisability.",
        "Short post-EV window; the transition is still unfolding.",
        "Reliance on reported figures; currency conventions (RMB vs. USD) differ across reports.",
        "NIO's capital-market model is noted but outside this team's scope.",
    ], 13, gap=9, marker_color=AMBER)

    lx = MARGIN + Inches(6.15)
    rect(s, lx, Inches(1.9), Inches(5.95), Inches(3.9), CARD, line_color=LINE)
    rect(s, lx, Inches(1.9), Inches(5.95), Inches(0.09), BLUE)
    text(s, lx + Inches(0.2), Inches(2.1), Inches(5.55), Inches(0.4),
         "Future research", 16, BLUE, font=BODY_FONT, bold=True)
    bullets(s, lx + Inches(0.2), Inches(2.6), Inches(5.55), Inches(3.0), [
        "Add a fourth archetype (NIO, capital-market model) for a fuller typology.",
        "Track outcomes over a longer EV cycle to test durability of the ranking.",
        "Quantify governance with a formal index and test it across more emerging markets.",
        "Examine how EU tariffs reshape each model's export strategy.",
    ], 13, gap=9, marker_color=BLUE)
    footer(s, 16)
    notes(s, """A short word on limits, since questions will follow. We study three firms in one
industry and one country, so we cannot claim a statistical law. The EV shift is still young, so
today's ranking may move. We rely on reported numbers, and the reports mix currencies, which we
had to reconcile. NIO would add a fourth model but was out of scope. Future work could add NIO,
follow the firms over a longer period, and build a formal governance score. Now a brief, honest
look at how we used AI in this project.""")


def s17_ai_platform(prs):
    s = blank(prs)
    kicker_headline(s, "AI annex · Platform & setup",
                    "AI used as research infrastructure — a layered agent architecture")
    # models
    text(s, MARGIN, Inches(1.8), Inches(12.1), Inches(0.3),
         "PLATFORM: CLAUDE CODE (ANTHROPIC)", 12, BLUE, font=BODY_FONT, bold=True)
    models = [("claude-sonnet-4-6", "Primary model: ingestion, drafting, APA checks, comparison table"),
              ("claude-opus-4", "Selective: cross-case reasoning, central argument, logic checks")]
    x = MARGIN
    for m, d in models:
        rect(s, x, Inches(2.15), Inches(5.95), Inches(1.0), CARD, line_color=LINE)
        rect(s, x, Inches(2.15), Inches(0.09), Inches(1.0), BLUE)
        text(s, x + Inches(0.2), Inches(2.28), Inches(5.55), Inches(0.35),
             m, 14, NAVY, font=BODY_FONT, bold=True)
        text(s, x + Inches(0.2), Inches(2.65), Inches(5.55), Inches(0.45),
             d, 12, INK, font=BODY_FONT, line_spacing=1.0)
        x += Inches(6.15)

    # agents
    text(s, MARGIN, Inches(3.4), Inches(12.1), Inches(0.3),
         "FIVE SPECIALISED AGENTS", 12, BLUE, font=BODY_FONT, bold=True)
    agents = [("wiki-ingest", "Structured\nsource ingestion"),
              ("wiki-query", "Retrieval from\nknowledge base"),
              ("wiki-lint", "Health-check\n& consistency"),
              ("verify-claims", "Trace claims to\nprimary source"),
              ("find-references", "Find supporting\nliterature")]
    aw = Inches(2.32)
    x = MARGIN
    for name, desc in agents:
        rect(s, x, Inches(3.75), aw, Inches(1.35), CARD, line_color=LINE)
        rect(s, x, Inches(3.75), aw, Inches(0.07), TEAL)
        text(s, x + Inches(0.12), Inches(3.9), aw - Inches(0.24), Inches(0.35),
             name, 12.5, TEAL, font=BODY_FONT, bold=True)
        text(s, x + Inches(0.12), Inches(4.3), aw - Inches(0.24), Inches(0.7),
             desc, 11.5, INK, font=BODY_FONT, line_spacing=1.0)
        x += Inches(2.42)

    bullets(s, MARGIN, Inches(5.4), Inches(12.1), Inches(1.3), [
        "Shared git repository: one wiki subfolder per case (SAIC / BYD / Geely) plus a shared comparison layer.",
        "Source-quality rating per source: journal tier Q1–Q5 × case relevance R1–R4.",
        "Cost: Claude Pro at ~€21 / month per member; no pay-per-query, no paid data services.",
    ], 13, gap=7)
    footer(s, 17)
    notes(s, """Now the AI annex, which the seminar asks us to share openly. We used Claude Code as a
research tool, not a writer. Two models did the work: a fast one for most tasks and a stronger
one for hard analysis. Instead of one chat, we used five small agents, each with a job, like
ingesting sources or checking claims against the original. Everything lived in a shared git
repo, one folder per company, with a quality rating on every source. The cost was low, about 21
euros a month each. Next, where it actually helped most.""")


def s18_ai_value(prs):
    s = blank(prs)
    kicker_headline(s, "AI annex · What helped most",
                    "Most valuable: catching three real data errors before submission")
    errs = [
        ("EU tariff", "Provisional 36.3% → final 35.3% for SAIC, traced to the regulation."),
        ("Currency unit", "A net-profit figure inflated 10× by a unit ambiguity in a secondary source."),
        ("Share count", "A BYD EPS figure built on the wrong share count, fixed via the audited report."),
    ]
    cw = Inches(3.93)
    x = MARGIN
    for i, (t, d) in enumerate(errs):
        rect(s, x, Inches(1.85), cw, Inches(1.95), CARD, line_color=LINE)
        rect(s, x, Inches(1.85), cw, Inches(0.09), GREEN)
        text(s, x + Inches(0.18), Inches(2.05), cw - Inches(0.36), Inches(0.4),
             f"Error {i+1} · {t}", 14, GREEN, font=BODY_FONT, bold=True)
        text(s, x + Inches(0.18), Inches(2.5), cw - Inches(0.36), Inches(1.2),
             d, 12.5, INK, font=BODY_FONT, line_spacing=1.05)
        x += Inches(4.08)

    text(s, MARGIN, Inches(4.05), Inches(12.1), Inches(0.3),
         "WHAT WORKED IN PRACTICE", 12, BLUE, font=BODY_FONT, bold=True)
    bullets(s, MARGIN, Inches(4.4), Inches(12.1), Inches(2.1), [
        "Best prompt pattern: “review this as a Dozent looking for overclaims, logical gaps and APA errors.”",
        "Fact-checking: isolate one claim and ask it to trace the exact source passage before it enters the draft.",
        "Cross-case work was reliable only with explicit wiki access — “compare X using only verified wiki data,” not general recall.",
        "Persistence mattered: after a hardware shutdown the full research base was recovered from the wiki and memory.",
    ], 13.5, gap=8)
    footer(s, 18)
    notes(s, """The single most useful thing was error catching. The tool found three real data
mistakes before they reached the paper: a wrong EU tariff, a profit figure inflated ten times by
a unit mix-up, and an EPS based on the wrong share count. All were fixed against audited reports.
On method, two habits worked best. Telling it to act like a strict professor gave sharper
feedback than asking for a friendly review. And forcing it to trace each fact to the source
caught the errors. The honest limits come last.""")


def s19_ai_limits(prs):
    s = blank(prs)
    kicker_headline(s, "AI annex · Limits & lessons learned",
                    "AI as research infrastructure, never as ghostwriter")
    rect(s, MARGIN, Inches(1.9), Inches(5.95), Inches(3.7), CARD, line_color=LINE)
    rect(s, MARGIN, Inches(1.9), Inches(5.95), Inches(0.09), AMBER)
    text(s, MARGIN + Inches(0.2), Inches(2.1), Inches(5.55), Inches(0.4),
         "Limits encountered", 16, AMBER, font=BODY_FONT, bold=True)
    bullets(s, MARGIN + Inches(0.2), Inches(2.6), Inches(5.55), Inches(2.9), [
        "No access to paywalled journals (e.g. AER behind the Goethe paywall).",
        "Geely PDFs were not machine-readable — manual data extraction needed.",
        "Several incorrect DOIs had to be checked and corrected by hand.",
        "Currency conventions (USD vs. RMB) drifted and required manual reconciliation.",
    ], 12.5, gap=8, marker_color=AMBER)

    lx = MARGIN + Inches(6.15)
    rect(s, lx, Inches(1.9), Inches(5.95), Inches(3.7), CARD, line_color=LINE)
    rect(s, lx, Inches(1.9), Inches(5.95), Inches(0.09), BLUE)
    text(s, lx + Inches(0.2), Inches(2.1), Inches(5.55), Inches(0.4),
         "Lessons learned", 16, BLUE, font=BODY_FONT, bold=True)
    bullets(s, lx + Inches(0.2), Inches(2.6), Inches(5.55), Inches(2.9), [
        "Learning curve: several hours to use the agent architecture productively, not as a chat tool.",
        "AI output always needs critical review — early drafts contained overclaims.",
        "All analysis, argument and interpretation stayed fully independent of the platform.",
        "Best framing: a source-quality enforcer and critical reviewer, not a writer.",
    ], 12.5, gap=8, marker_color=BLUE)

    rect(s, MARGIN, Inches(5.85), Inches(12.1), Inches(0.7), NAVY)
    text(s, MARGIN + Inches(0.2), Inches(5.85), Inches(11.7), Inches(0.7),
         "Bottom line: AI compressed research time and caught errors — but every judgment, "
         "argument and conclusion remained the authors' own.",
         13.5, WHITE, font=BODY_FONT, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    footer(s, 19)
    notes(s, """Now the limits, stated plainly. The tool could not open paywalled journals, could not
read Geely's PDFs, invented some wrong DOIs, and sometimes mixed currencies. We fixed all of
this by hand. The biggest lesson was that it took real time to learn — it is not a chat box, and
its output always needs checking. Early drafts had overclaims we removed. The key principle: we
used AI as a fact-checker and editor, never as a ghostwriter. Every argument is ours. Finally,
our references.""")


def s20_references(prs):
    s = blank(prs)
    kicker_headline(s, "References", "Key sources (APA 7)")
    refs_l = [
        "Bai, J., Barwick, P. J., Cao, S., & Li, S. (2025). Quid pro quo, knowledge spillovers, and industrial quality upgrading. American Economic Review, 115(11).",
        "Balcet, G., Wang, H., & Richet, X. (2012). Geely: A trajectory of catching up and asset-seeking multinational growth. IJATM, 12(4).",
        "European Commission. (2024). Implementing Regulation (EU) 2024/2754. Official Journal of the EU.",
        "Fan, J. P. H., Wong, T. J., & Zhang, T. (2007). Politically connected CEOs and post-IPO performance. JFE, 84(2).",
        "Holmes, T. J., McGrattan, E. R., & Prescott, E. C. (2015). Quid pro quo: Technology capital transfers for market access in China. RES, 82(3).",
        "Howell, S. T. (2018). Joint ventures and technology adoption: A Chinese industrial policy that backfired. Research Policy, 47(8).",
        "Jensen, M. C., & Meckling, W. H. (1976). Theory of the firm. JFE, 3(4).",
        "Konda, P., Slepnikov, D., & Jin, J. (2022). From transaction to co-creation in Geely's acquisition of Volvo Cars. AJTI, 31(3).",
    ]
    refs_r = [
        "Luo, Y., & Tung, R. L. (2007). International expansion of emerging market enterprises: A springboard perspective. JIBS, 38(4).",
        "Megginson, W. L., & Netter, J. M. (2001). From state to market: A survey of empirical studies on privatization. JEL, 39(2).",
        "Shleifer, A., & Vishny, R. W. (1994). Politicians and firms. Quarterly Journal of Economics, 109(4).",
        "Teece, D. J., Pisano, G., & Shuen, A. (1997). Dynamic capabilities and strategic management. SMJ, 18(7).",
        "Whitfield, L., & Wuttke, T. (2026). China's technological catch-up and leapfrogging in electric vehicles. Progress in Economic Geography, 4.",
        "International Energy Agency. (2026). Global EV outlook 2026. IEA.",
        "Annual Reports: SAIC Motor (2025, 2026); BYD (2025, 2026); Geely Automobile Holdings (2025a, 2025b, 2026).",
        "Also cited: Bortolotti & Faccio (2009); Inkpen & Beamish (1997); Wang et al. (2022); Zhang et al. (2025); Zheng et al. (2022); Allen et al. (2024).",
    ]
    bullets(s, MARGIN, Inches(1.85), Inches(5.95), Inches(5.0), refs_l, 10.5,
            gap=8, marker="", line_spacing=1.0)
    bullets(s, MARGIN + Inches(6.15), Inches(1.85), Inches(5.95), Inches(5.0), refs_r, 10.5,
            gap=8, marker="", line_spacing=1.0)
    footer(s, 20)
    notes(s, """These are the main sources behind the talk, in APA 7 style. The core theory comes from
Shleifer and Vishny, Jensen and Meckling, Teece, Luo and Tung, and Fan. The key empirical work
on China is Bai, Howell, and Holmes. Company figures come from the audited annual reports of
SAIC, BYD, and Geely, plus the EU regulation and the IEA outlook. The full list is in the paper.
Thank you — we are happy to take questions.""")


# ═══════════════════════════════════════════════════════════════════════════════

def main():
    out_dir = r"C:\Users\User\CHINA-seminar-paper\research\presentation"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "china-strategies-julia.pptx")

    prs = new_prs()
    for fn in [s01_title, s02_context, s03_question, s04_theory, s05_framework,
               s06_governance, s07_governance_synth, s08_financing, s09_financing_synth,
               s10_capability, s11_capability_synth, s12_viability, s13_scorecard,
               s14_lessons12, s15_lesson3_concl, s16_limitations,
               s17_ai_platform, s18_ai_value, s19_ai_limits, s20_references]:
        fn(prs)

    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
