"""
McKinsey-style slide generator — editable PPTX via python-pptx.
Usage: uv run generate_slides.py
Output: Slides/output.pptx
"""

# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx"]
# ///

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x00, 0x2D, 0x62)   # McKinsey navy
BLUE_MID    = RGBColor(0x00, 0x58, 0xA3)   # section headers
BLUE_ACCENT = RGBColor(0x00, 0xA3, 0xE0)   # accent / underline
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x63, 0x63, 0x63)
LIGHT_GRAY  = RGBColor(0xF4, 0xF6, 0xF9)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
RED         = RGBColor(0xC4, 0x1E, 0x3A)
GREEN       = RGBColor(0x1B, 0x6B, 0x3A)

# ── Slide dimensions (16:9, 1280×720 px equivalent) ──────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(14), bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT, wrap=True,
             font_name="Arial"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def top_bar(slide, color=BLUE_DARK, height=Inches(0.08)):
    add_rect(slide, 0, 0, W, height, fill_rgb=color)


def page_number(slide, number):
    add_text(slide, str(number),
             left=W - Inches(0.6), top=H - Inches(0.35),
             width=Inches(0.5), height=Inches(0.3),
             font_size=Pt(11), color=GRAY, align=PP_ALIGN.RIGHT)


# ── Slide builders ────────────────────────────────────────────────────────────

def title_slide(prs, title, subtitle, authors, course_info, date):
    """Full dark-blue title slide."""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_rgb=BLUE_DARK)
    # accent bar
    add_rect(slide, 0, 0, W, Inches(0.1), fill_rgb=BLUE_ACCENT)

    # Title
    add_text(slide, title,
             left=Inches(0.9), top=Inches(1.8),
             width=Inches(11.5), height=Inches(1.4),
             font_size=Pt(40), bold=True, color=WHITE)
    # Accent line under title
    add_rect(slide, Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.04),
             fill_rgb=RGBColor(0xFF, 0xFF, 0xFF))

    # Subtitle
    add_text(slide, subtitle,
             left=Inches(0.9), top=Inches(3.5),
             width=Inches(11.5), height=Inches(0.6),
             font_size=Pt(20), color=RGBColor(0xCC, 0xD9, 0xEC))

    # Course info
    add_text(slide, course_info,
             left=Inches(0.9), top=Inches(5.0),
             width=Inches(11.5), height=Inches(0.5),
             font_size=Pt(14), color=RGBColor(0xA0, 0xB8, 0xD8))

    # Authors + date
    add_text(slide, f"{authors}  |  {date}",
             left=Inches(0.9), top=Inches(5.6),
             width=Inches(11.5), height=Inches(0.4),
             font_size=Pt(13), color=RGBColor(0x88, 0xA4, 0xC8))


def divider_slide(prs, section_number, section_title, subtitle="", page_num=None):
    """Blue section-divider slide."""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_rgb=BLUE_MID)
    add_rect(slide, 0, 0, W, Inches(0.08), fill_rgb=BLUE_ACCENT)

    # Section number
    add_text(slide, section_number,
             left=Inches(0.9), top=Inches(2.2),
             width=Inches(2), height=Inches(0.5),
             font_size=Pt(14), color=BLUE_ACCENT, bold=True)

    # Section title
    add_text(slide, section_title,
             left=Inches(0.9), top=Inches(2.7),
             width=Inches(11.5), height=Inches(1.2),
             font_size=Pt(38), bold=True, color=WHITE)

    # Divider line
    add_rect(slide, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.04),
             fill_rgb=RGBColor(0xFF, 0xFF, 0xFF))

    if subtitle:
        add_text(slide, subtitle,
                 left=Inches(0.9), top=Inches(4.2),
                 width=Inches(11.5), height=Inches(0.6),
                 font_size=Pt(18), color=RGBColor(0xCC, 0xD9, 0xEC))

    if page_num:
        page_number(slide, page_num)


def content_slide(prs, headline, key_message, bullets, page_num=None):
    """
    Standard content slide.
    headline    : slide title (str)
    key_message : bold insight line shown below headline (str or None)
    bullets     : list of strings. Prefix with "  " for sub-bullet.
    """
    slide = prs.slides.add_slide(blank_layout(prs))
    top_bar(slide)

    # Headline
    add_text(slide, headline,
             left=Inches(0.7), top=Inches(0.2),
             width=Inches(11.9), height=Inches(0.75),
             font_size=Pt(24), bold=True, color=BLUE_DARK)
    # Accent line under headline
    add_rect(slide, Inches(0.7), Inches(0.98), Inches(11.9), Inches(0.03),
             fill_rgb=BLUE_ACCENT)

    y = Inches(1.15)

    # Key message box
    if key_message:
        box = add_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.55),
                       fill_rgb=BLUE_DARK)
        add_text(slide, key_message,
                 left=Inches(0.85), top=y + Inches(0.06),
                 width=Inches(11.6), height=Inches(0.43),
                 font_size=Pt(14), bold=True, color=WHITE)
        y += Inches(0.65)

    # Bullets
    for item in bullets:
        is_sub = item.startswith("  ")
        text   = item.lstrip()
        indent = Inches(1.1) if is_sub else Inches(0.7)
        bullet = "–" if is_sub else "▪"
        color  = GRAY if is_sub else BLACK
        fsize  = Pt(13) if is_sub else Pt(14.5)

        # bullet symbol
        add_text(slide, bullet,
                 left=indent, top=y,
                 width=Inches(0.25), height=Inches(0.32),
                 font_size=fsize, color=BLUE_ACCENT if not is_sub else GRAY)
        # text
        add_text(slide, text,
                 left=indent + Inches(0.25), top=y,
                 width=Inches(11.6) - indent, height=Inches(0.32),
                 font_size=fsize, color=color)
        y += Inches(0.36)

    if page_num:
        page_number(slide, page_num)


def table_slide(prs, headline, key_message, headers, rows, page_num=None,
                col_widths=None):
    """
    Slide with a McKinsey-style table.
    headers    : list of column header strings
    rows       : list of lists (each inner list = one row of cells)
    col_widths : list of Inches values (must sum to ~11.9); auto if None
    """
    slide = prs.slides.add_slide(blank_layout(prs))
    top_bar(slide)

    add_text(slide, headline,
             left=Inches(0.7), top=Inches(0.2),
             width=Inches(11.9), height=Inches(0.75),
             font_size=Pt(24), bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(0.7), Inches(0.98), Inches(11.9), Inches(0.03),
             fill_rgb=BLUE_ACCENT)

    y = Inches(1.15)

    if key_message:
        add_rect(slide, Inches(0.7), y, Inches(11.9), Inches(0.55),
                 fill_rgb=BLUE_DARK)
        add_text(slide, key_message,
                 left=Inches(0.85), top=y + Inches(0.06),
                 width=Inches(11.6), height=Inches(0.43),
                 font_size=Pt(14), bold=True, color=WHITE)
        y += Inches(0.65)

    # Build table
    n_cols = len(headers)
    n_rows = len(rows) + 1   # +1 for header row

    if col_widths is None:
        cw = Inches(11.9) / n_cols
        col_widths = [cw] * n_cols

    table_h = Inches(0.42) * n_rows
    tbl = slide.shapes.add_table(n_rows, n_cols,
                                 Inches(0.7), y,
                                 Inches(11.9), table_h).table

    # Column widths
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    def style_cell(cell, text, bg, fg, bold=False, font_size=Pt(13)):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name  = "Arial"
        run.font.size  = font_size
        run.font.bold  = bold
        run.font.color.rgb = fg

    # Header row
    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, BLUE_DARK, WHITE, bold=True, font_size=Pt(13))

    # Data rows
    for i, row in enumerate(rows):
        bg = LIGHT_GRAY if i % 2 == 1 else WHITE
        for j, val in enumerate(row):
            style_cell(tbl.cell(i + 1, j), str(val), bg, BLACK, font_size=Pt(12))

    if page_num:
        page_number(slide, page_num)


def quote_slide(prs, headline, quote, attribution="", bullets=None, page_num=None):
    """Slide with a prominent pull-quote box plus optional bullets."""
    slide = prs.slides.add_slide(blank_layout(prs))
    top_bar(slide)

    add_text(slide, headline,
             left=Inches(0.7), top=Inches(0.2),
             width=Inches(11.9), height=Inches(0.75),
             font_size=Pt(24), bold=True, color=BLUE_DARK)
    add_rect(slide, Inches(0.7), Inches(0.98), Inches(11.9), Inches(0.03),
             fill_rgb=BLUE_ACCENT)

    # Quote box
    add_rect(slide, Inches(0.7), Inches(1.15), Inches(11.9), Inches(1.0),
             fill_rgb=LIGHT_GRAY)
    add_rect(slide, Inches(0.7), Inches(1.15), Inches(0.07), Inches(1.0),
             fill_rgb=BLUE_ACCENT)
    add_text(slide, f'"{quote}"',
             left=Inches(0.95), top=Inches(1.22),
             width=Inches(11.4), height=Inches(0.7),
             font_size=Pt(16), italic=True, color=BLUE_DARK)
    if attribution:
        add_text(slide, f"— {attribution}",
                 left=Inches(0.95), top=Inches(1.9),
                 width=Inches(11.4), height=Inches(0.28),
                 font_size=Pt(12), color=GRAY)

    y = Inches(2.35)
    if bullets:
        for item in bullets:
            is_sub = item.startswith("  ")
            text   = item.lstrip()
            indent = Inches(1.1) if is_sub else Inches(0.7)
            bullet = "–" if is_sub else "▪"
            color  = GRAY if is_sub else BLACK
            fsize  = Pt(13) if is_sub else Pt(14.5)
            add_text(slide, bullet,
                     left=indent, top=y,
                     width=Inches(0.25), height=Inches(0.32),
                     font_size=fsize, color=BLUE_ACCENT if not is_sub else GRAY)
            add_text(slide, text,
                     left=indent + Inches(0.25), top=y,
                     width=Inches(11.6) - indent, height=Inches(0.32),
                     font_size=fsize, color=color)
            y += Inches(0.36)

    if page_num:
        page_number(slide, page_num)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE CONTENT — edit everything below this line
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = new_prs()

    # ── Slide 1 — Title ───────────────────────────────────────────────────────
    title_slide(prs,
        title="China Strategies in the Automotive Industry",
        subtitle="Governance, Financing, and Long-Term Value Creation",
        authors="Julia Ort · Batughan Maus · Chiara Stubner",
        course_info='Seminar "M&A as Transformer of the Automotive Industry" | Goethe-Universität Frankfurt',
        date="19. Juni 2026"
    )

    # ── Slide 2 — Section divider: Introduction ───────────────────────────────
    divider_slide(prs,
        section_number="01",
        section_title="Introduction",
        subtitle="Setting the Stage",
        page_num=2
    )

    # ── Slide 3 — 55% EV share ────────────────────────────────────────────────
    content_slide(prs,
        headline="55% of all new cars in China are electric in 2025 — who survives?",
        key_message="The market has already tipped — incumbents must adapt or exit",
        bullets=[
            "EV share China 2025: ~55% of all new registrations (IEA 2026)",
            "Chinese OEMs: 60% of all global EV sales",
            "70% of BEVs sold in China cost less than the average car",
            "IEA forecast: >90% EV share in China by 2035",
            "",
            "Research Question: Which governance and financing structure creates long-term",
            "  shareholder value in a technology-intensive, fast-moving market?",
            "",
            "Three cases under identical market conditions: SAIC · BYD · Geely",
        ],
        page_num=3
    )

    # ── Slide 4 — Three models overview ──────────────────────────────────────
    table_slide(prs,
        headline="Three paths, one market — sharply divergent outcomes",
        key_message="Governance structure drives divergence in EV readiness and financial performance",
        headers=["", "SAIC Motor", "BYD", "Geely"],
        rows=[
            ["Model",          "State-directed JV",         "Organic vertical integration",  "M&A-driven private"],
            ["Control",        "Shanghai SASAC (63.71%)",   "Wang Chuanfu, founder (17%)",   "Li Shufu, private (100% ZGH)"],
            ["Net profit 2024","RMB 1.67bn (−88%)",         "RMB 40.25bn (+34%)",            "RMB 16.8bn (+240%)"],
            ["EV Readiness",   "Medium — JV restructuring", "High — ICE exit 2022",          "Medium — platform-dependent"],
        ],
        col_widths=[Inches(2.0), Inches(3.3), Inches(3.3), Inches(3.3)],
        page_num=4
    )

    # ── Slide 5 — Section divider: Theory ────────────────────────────────────
    divider_slide(prs,
        section_number="02",
        section_title="Theory & Literature Review",
        subtitle="Four Theoretical Lenses",
        page_num=5
    )

    # ── Slide 6 — Theory ─────────────────────────────────────────────────────
    content_slide(prs,
        headline="Four theoretical lenses — one analytical framework",
        key_message=None,
        bullets=[
            "1. Political Governance & SOE Efficiency",
            "  Shleifer & Vishny (1994): politicians maximize employment, not firm value",
            "  Fan, Wong & Zhang (2007): politically connected CEOs → −18% CAR, −21% revenue growth",
            "  Megginson & Netter (2001): mixed-ownership firms no more profitable than pure SOEs",
            "",
            "2. JV Dynamics & Capability Transfer",
            "  Holmes, McGrattan & Prescott (2015): JV as implicit technology tax",
            '  Howell (2018): JV rents crowd out autonomous innovation ("like opium")',
            "",
            "3. Cross-border M&A & Springboard Strategy",
            "  Luo & Tung (2007): emerging-market firms use acquisitions as springboard",
            "  Zheng, Noorderhaven & Du (2022): 4-stage integration (distancing → diversifying)",
            "",
            "4. Dynamic Capabilities & Founder Control",
            "  Teece, Pisano & Shuen (1997): sustainable advantage through path-dependent tacit knowledge",
            "  Zhang et al. (2025): founder control → more breakthrough innovation",
        ],
        page_num=6
    )

    # ── Slide 7 — Section divider: Governance ────────────────────────────────
    divider_slide(prs,
        section_number="03",
        section_title="Governance",
        subtitle="Who Controls the Steering Wheel?",
        page_num=7
    )

    # ── Slide 8 — SAIC Governance ─────────────────────────────────────────────
    table_slide(prs,
        headline="SAIC: State as controller — politically embedded objectives",
        key_message="~73% state-owned → political goals embedded in governance structure",
        headers=["Metric", "Politically connected", "Non-connected"],
        rows=[
            ["3-year CAR",       "−18%",  "Baseline"],
            ["Revenue growth",   "−21%",  "Baseline"],
            ["Profit growth",    "−24%",  "Baseline"],
            ["Board bureaucrats","36%",   "19%"],
        ],
        col_widths=[Inches(4.5), Inches(3.7), Inches(3.7)],
        page_num=8
    )

    # ── Slide 9 — BYD & Geely Governance ─────────────────────────────────────
    content_slide(prs,
        headline="BYD & Geely: Founder control without political constraints",
        key_message="Consolidated control enables fast, unilateral strategic decisions",
        bullets=[
            "BYD — Wang Chuanfu: Chairman + CEO since 1995, ~17% direct; >1/3 voting rights",
            "  No Party Committee directive; subsidies (RMB 10.4bn) grant no control rights",
            "  ICE exit 2022, BYD Semiconductor, FinDreams spin-off — no external approval needed",
            '  Wang Chuanfu: "Fast decision-making" as central competitive advantage (He 2025)',
            "",
            "Geely — Li Shufu: 100% owner ZGH (unlisted apex) → ~76% control over Geely Auto HK",
            "  No JV-partner or state-committee approval needed",
            "  Volvo 2010 ($1.8bn) → LEVC 2013 → Lotus 2017 (51%) → Smart JV 2020 (50:50 Mercedes)",
            "  Risk: multi-brand portfolio generates growing coordination costs",
        ],
        page_num=9
    )

    # ── Slide 10 — Section divider: Financing ────────────────────────────────
    divider_slide(prs,
        section_number="04",
        section_title="Financing",
        subtitle="Capital Structure as Strategic Bet",
        page_num=10
    )

    # ── Slide 11 — SAIC & BYD Financing ──────────────────────────────────────
    content_slide(prs,
        headline="SAIC & BYD: Capital trap vs. self-financing resilience",
        key_message="SAIC absorbs 50% impairment without restructuring control; BYD reinvests all cash flow",
        bullets=[
            "SAIC — Capital-light JV model becomes a capital trap",
            "  2024: SAIC-GM sales −56.54% → must absorb 50% impairment without control",
            "  Revenue −15.73% (RMB 627.59bn), Net profit −88.19% (RMB 1.67bn), ROE 0.58%",
            "  EU countervailing duty: 35.3% — highest of all three cases",
            "",
            "BYD — Internal financing as structural strength",
            "  FY2024: Revenue RMB 777.1bn (+29%), Operating CF RMB 133.5bn, Cash RMB 154.9bn",
            "  R&D 2024: RMB 54.2bn (+35.7%) — exceeds annual net profit RMB 40.3bn",
            "  In 13 of 14 years (2011–2024): R&D expenditure exceeded annual net profit (He 2025)",
            "  Interest-bearing debt only 4.9% of total liabilities",
        ],
        page_num=11
    )

    # ── Slide 12 — Geely Financing ────────────────────────────────────────────
    content_slide(prs,
        headline="Geely: Leverage for speed — highest financial risk",
        key_message="Volvo deal returned 10× — but Polestar and Lotus now create portfolio liability",
        bullets=[
            "Volvo Deal 2010 — The paradigm case",
            "  Purchase price: $1.8bn → total capital ~$2.7bn",
            "  Financing: Bank of China/CCB/ExIm ~$2.1bn + ~$932m central government",
            "  ZGH debt: RMB 16.05bn (2009) → RMB 71.07bn (2010); D/C ratio 73.47%",
            "  Return: Volvo IPO October 2021 ~$18.5bn — 10× purchase price",
            "",
            "Current risks (2024/25)",
            "  Polestar: Net loss $1.17bn FY2023, going-concern warning (Deloitte), $103m cash end-2024",
            "  Polestar: ZGH as guarantor, ~$2bn capital need through 2028",
            "  Lotus Technology: cumulative deficit ~$3.2bn (FY2025), Revenue −44% FY2025",
        ],
        page_num=12
    )

    # ── Slide 13 — Section divider: Capability Building ──────────────────────
    divider_slide(prs,
        section_number="05",
        section_title="Capability Building",
        subtitle="Build, Buy, or Borrow?",
        page_num=13
    )

    # ── Slide 14 — SAIC Capability Paradox ───────────────────────────────────
    quote_slide(prs,
        headline="SAIC and the Capability Paradox",
        quote="JV rents are like opium",
        attribution="He Guangyan, former Minister of Machinery",
        bullets=[
            "What the JV model gave (Bai et al. 2025, n=all auto models 2001–2014)",
            "  +8.3% quality improvement in JV-affiliated models vs. non-affiliated",
            "  Defect rate: 65% higher in 2001 → 33% higher in 2014",
            "  Transfer via worker mobility + supplier networks (not patent transfer)",
            "",
            "What the JV model prevented (Howell 2018)",
            "  JV firms responded to regulatory pressure with quality reduction, not innovation",
            "  Cannibalization: own brand competes with profitable JV revenue stream",
            "",
            "Capability Paradox: JV transferred production quality but suppressed autonomous",
            "  innovation needed for the EV era — no proprietary battery, no software platform",
        ],
        page_num=14
    )

    # ── Slide 15 — BYD & Geely Capability ────────────────────────────────────
    content_slide(prs,
        headline="BYD & Geely: Organic depth vs. M&A speed",
        key_message="BYD's tacit knowledge is non-imitable; Geely's M&A created genuine co-development",
        bullets=[
            "BYD — Organic internal build: structurally hard to imitate",
            "  1995: mobile phone batteries → 2003: vehicles → mid-2000s: automotive batteries",
            '  "Specialised Vertical Integration": competencies stack on each other (Wang, Zhao & Ruet 2022)',
            "  Teece et al. (1997): tacit knowledge in routines — not transferable via licence or talent",
            "  BYD did not follow policy agenda — it actively co-designed it (Whitfield & Wuttke 2026)",
            "",
            "Geely — M&A as capability compressor",
            "  Volvo 2010: immediate access to safety engineering, 3 platforms, 2,400+ dealer network",
            "  4-stage integration (Zheng et al. 2022): distancing → balancing → building → diversifying",
            "  CMA Platform: bilateral co-development across Volvo, Lynk & Co, Geely",
            "  R&D expenditure: +600% in 7 years post-acquisition",
            "  Konda et al. (2022): measurable patent convergence → genuine co-creation effect",
        ],
        page_num=15
    )

    # ── Slide 16 — Section divider: Long-Term Viability ──────────────────────
    divider_slide(prs,
        section_number="06",
        section_title="Long-Term Viability",
        subtitle="Who Has a Sustainable Position?",
        page_num=16
    )

    # ── Slide 17 — SAIC Viability ─────────────────────────────────────────────
    table_slide(prs,
        headline="SAIC: JV erosion meets EV transition simultaneously",
        key_message="Double pressure: JV demand collapsing while EV transition requires new capital allocation",
        headers=["Facility", "Capacity", "Production", "Utilisation"],
        rows=[
            ["SAIC-GM",       "1.452m",    "534,000",   "37%"],
            ["SAIC-VW",       "1.920m",    "1.058m",    "55%"],
            ["VW MEB (EV)",   "240,000",   "34,464",    "14%"],
            ["Own-brand (PVB)","—",         "—",         "76–95%"],
        ],
        col_widths=[Inches(3.5), Inches(2.8), Inches(2.8), Inches(2.8)],
        page_num=17
    )

    # ── Slide 18 — BYD & Geely Viability ─────────────────────────────────────
    table_slide(prs,
        headline="BYD & Geely: EV transition complete vs. strong core, risky periphery",
        key_message="BYD structurally repositioned; Geely strong at core but portfolio liabilities accumulate",
        headers=["Entity", "Performance 2024/25"],
        rows=[
            ["BYD",             ">1m exported vehicles (+140% YoY), 119 countries, 4th year global NEV leader"],
            ["BYD risk",        "Gross margin 19%→18%; Net profit −19% (RMB 32.6bn); net-cash → net-debt"],
            ["Geely Auto",      "Revenue +32% ($33.4bn), Net profit +240%"],
            ["Volvo Cars",      "SEK 400.2bn revenue, 6.8% EBIT — second record year"],
            ["Geely stock",     "+39% early 2026 vs. BYD +12%"],
            ["Polestar",        "Going concern, $103m cash, ~$2bn capital need, ZGH guarantor"],
            ["Lotus Technology","Revenue −44% FY2025, cumulative deficit ~$3.2bn"],
        ],
        col_widths=[Inches(2.5), Inches(9.4)],
        page_num=18
    )

    # ── Slide 19 — Section divider: Lessons ──────────────────────────────────
    divider_slide(prs,
        section_number="07",
        section_title="Generalizable Lessons",
        subtitle="What Can Other Markets Learn?",
        page_num=19
    )

    # ── Slide 20 — Lesson 1 ───────────────────────────────────────────────────
    table_slide(prs,
        headline="Lesson 1: Mode of capability acquisition determines durability across technology cycles",
        key_message="Mandatory Technology Transfer creates short-term gains but crowds out next-cycle innovation",
        headers=["", "SAIC (Mandatory JV)", "BYD (Organic)", "Geely (M&A)"],
        rows=[
            ["Short-term", "+8.3% quality (Bai et al.)", "Slow",                   "Immediate platforms"],
            ["Long-term",  "Autonomous innovation crowded out", "Proprietary, non-imitable", "Bilateral co-development"],
            ["EV readiness","Missing",                   "High",                  "Platform-dependent"],
            ["Cost",       "Political control",          "Time",                  "Leverage & complexity"],
        ],
        col_widths=[Inches(1.8), Inches(3.37), Inches(3.37), Inches(3.36)],
        page_num=20
    )

    # ── Slide 21 — Lessons 2 & 3 ─────────────────────────────────────────────
    quote_slide(prs,
        headline="Lessons 2 & 3: Governance and capital structure as first-order constraints",
        quote="When the rate of technological change exceeds the decision cycle time of a governance structure, governance itself becomes the binding competitive constraint.",
        attribution=None,
        bullets=[
            "Lesson 2: Governance is a first-order determinant of strategic adaptability",
            "  SAIC: control split across SASAC + Party Committee + 50:50 JV boards → slow",
            "  BYD: consolidated founder control → ICE exit 2022 without external approval",
            "  Geely: private apex, but multi-brand portfolio generates coordination costs",
            "",
            "Lesson 3: Capital structure determines resilience when markets contract",
            "  SAIC: capital-light → capital trap; 50% impairment without restructuring control",
            "  BYD: self-financing → resilient to partner losses; vulnerable to margin compression",
            "  Geely: leverage for speed → 10× Volvo return; highest portfolio liability risk",
        ],
        page_num=21
    )

    # ── Slide 22 — Section divider: Conclusions ──────────────────────────────
    divider_slide(prs,
        section_number="08",
        section_title="Conclusions",
        subtitle="Governance as Binding Constraint",
        page_num=22
    )

    # ── Slide 23 — Conclusions ────────────────────────────────────────────────
    table_slide(prs,
        headline="Governance is the binding constraint across all four dimensions",
        key_message="Governance determines financing options, capability development, and competitive position",
        headers=["", "Governance", "Financing", "Capabilities", "EV Readiness", "Verdict"],
        rows=[
            ["BYD",   "Founder ✓",   "Self-fin. ✓", "Organic ✓",   "Complete ✓",      "Structurally superior"],
            ["Geely", "Private ✓",   "Leverage ⚠",  "M&A ✓",       "Platform-dep. ⚠", "Strong core, risky periphery"],
            ["SAIC",  "State/Party ✗","Capital trap ✗","JV-dep. ✗","Transitioning ✗", "Double structural brake"],
        ],
        col_widths=[Inches(1.0), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(3.7)],
        page_num=23
    )

    # ── Slide 24 — AI Highlights ──────────────────────────────────────────────
    content_slide(prs,
        headline="AI Highlights — Shared Git Repository + Claude (Anthropic)",
        key_message=None,
        bullets=[
            "What we used",
            "  Shared git repository with three separate company wikis (BYD, SAIC, Geely)",
            "  Claude Code CLI with wiki-ingest, wiki-query, wiki-lint agents",
            "  29+ academic sources structured, ingested, cross-referenced",
            "  Source-rating system for source prioritisation",
            "",
            "What worked well",
            "  Wiki-Ingest: automatic extraction and linking of entities and concepts across sources",
            "  Cross-source synthesis: contradictions between sources automatically flagged",
            "  Writing assistance: sections drafted directly from cited sources",
            "",
            "Challenges",
            "  No direct SharePoint/cloud access → files had to be stored locally",
            "  Hallucinations on financial figures → consistent cross-checks against annual reports",
            "  Wiki schema required initial setup (~1–2h)",
            "",
            "Cost: Claude Pro ~$20/month",
        ],
        page_num=24
    )

    return prs


if __name__ == "__main__":
    prs = build_presentation()
    out = "Slides/output.pptx"
    prs.save(out)
    print(f"Saved: {out}")
