# -*- coding: utf-8 -*-
"""Build the China Strategies seminar presentation (20 slides) with speaker notes."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ---- palette ----
NAVY  = RGBColor(0x1F, 0x38, 0x64)
BLUE  = RGBColor(0x2E, 0x54, 0x96)
ACC   = RGBColor(0xED, 0x7D, 0x31)   # orange accent
LIGHT = RGBColor(0xED, 0xF1, 0xF7)
GRAY  = RGBColor(0x3A, 0x3A, 0x3A)
MGRAY = RGBColor(0x7A, 0x7A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def _set(font, size=None, bold=None, color=None, name="Calibri", italic=None):
    if size  is not None: font.size  = Pt(size)
    if bold  is not None: font.bold  = bold
    if color is not None: font.color.rgb = color
    if italic is not None: font.italic = italic
    font.name = name

def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def footer(slide, n):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(9), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "China Strategies in the Automotive Industry"
    _set(r.font, 9, False, MGRAY)
    tb2 = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = str(n); _set(r2.font, 9, False, MGRAY)

def header(slide, section, title):
    bar = rect(slide, 0, 0, SW, Inches(1.25), NAVY)
    # section label
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(12), Inches(0.35))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = section.upper()
    _set(r.font, 12, True, ACC)
    # title
    tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(0.52), Inches(12.2), Inches(0.7))
    r2 = tb2.text_frame.paragraphs[0].add_run(); r2.text = title
    _set(r2.font, 24, True, WHITE)

def content(n, section, title, bullets, note):
    """bullets: list of (level, text) ; text may contain 'Lead: rest' for bold lead."""
    s = prs.slides.add_slide(BLANK)
    header(s, section, title)
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12.1), Inches(5.2))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8); p.space_before = Pt(2)
        # bold lead before ':' if short
        lead, sep, rest = text.partition(": ")
        bullet_char = "▪  " if level == 0 else "–  "
        if sep and len(lead) <= 26:
            r0 = p.add_run(); r0.text = bullet_char + lead + ": "
            _set(r0.font, 17 if level==0 else 15, True, NAVY if level==0 else BLUE)
            r1 = p.add_run(); r1.text = rest
            _set(r1.font, 17 if level==0 else 15, False, GRAY)
        else:
            r = p.add_run(); r.text = bullet_char + text
            _set(r.font, 17 if level==0 else 15, level==0, GRAY)
    footer(s, n)
    notes(s, note)
    return s

def style_cell(cell, text, size, bold, color, fill, align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    _set(r.font, size, bold, color)

def table_slide(n, section, title, data, note, col_widths=None, first_col_left=True, fontsize=11):
    s = prs.slides.add_slide(BLANK)
    header(s, section, title)
    rows, cols = len(data), len(data[0])
    tw = Inches(12.1); th = Inches(0.55*rows)
    gx = Inches(0.6); gy = Inches(1.7)
    gtbl = s.shapes.add_table(rows, cols, gx, gy, tw, th)
    tbl = gtbl.table
    # remove default styling band
    tbl.first_row = False; tbl.horz_banding = False
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            is_head = (ri == 0)
            is_label = (ci == 0)
            if is_head:
                fill, color, bold, sz = NAVY, WHITE, True, fontsize
            elif is_label:
                fill, color, bold, sz = BLUE, WHITE, True, fontsize
            else:
                fill = LIGHT if ri % 2 == 0 else WHITE
                color, bold, sz = GRAY, False, fontsize
            align = PP_ALIGN.LEFT if (is_label and first_col_left) or (is_head and ci==0) else PP_ALIGN.CENTER
            style_cell(tbl.cell(ri, ci), str(val), sz, bold, color, fill, align)
    footer(s, n)
    notes(s, note)
    return s

def stat_slide(n, section, title, stats, sub, note):
    s = prs.slides.add_slide(BLANK)
    header(s, section, title)
    n3 = len(stats)
    boxw = Inches(3.6); gap = Inches(0.55)
    total = boxw*n3 + gap*(n3-1)
    startx = (SW - total)//2
    for i,(num,label,col) in enumerate(stats):
        x = startx + i*(boxw+gap)
        card = rect(s, x, Inches(2.4), boxw, Inches(2.4), LIGHT)
        accentbar = rect(s, x, Inches(2.4), boxw, Inches(0.12), col)
        tb = s.shapes.add_textbox(x, Inches(2.85), boxw, Inches(1.2))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; _set(r.font, 48, True, col)
        tb2 = s.shapes.add_textbox(x, Inches(4.05), boxw, Inches(0.7))
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = label; _set(r2.font, 16, True, NAVY)
    if sub:
        tb3 = s.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0))
        tb3.text_frame.word_wrap = True
        p3 = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = sub; _set(r3.font, 16, False, GRAY)
    footer(s, n)
    notes(s, note)
    return s

# ============================ SLIDE 1 — TITLE ============================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.06), ACC)
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "China Strategies in the Automotive Industry"
_set(r.font, 40, True, WHITE)
p2 = tf.add_paragraph(); p2.space_before = Pt(14)
r2 = p2.add_run(); r2.text = "Which governance and financing model creates long-term shareholder value?"
_set(r2.font, 20, False, RGBColor(0xCB,0xD6,0xE6), italic=True)
tb2 = s.shapes.add_textbox(Inches(0.9), Inches(4.8), Inches(11.5), Inches(2.0))
tf2 = tb2.text_frame; tf2.word_wrap = True
for txt, sz, col in [
    ("Cluster: China strategies  ·  Seminar “M&A as Transformer of the Automotive Industry”", 15, RGBColor(0x9D,0xB4,0xD4)),
    ("Julia Ort  ·  Batughan Maus  ·  Chiara Stubner", 17, WHITE),
    ("Chair of Banking and Finance  ·  Prof. Dr. Mark Wahrenburg", 14, RGBColor(0x9D,0xB4,0xD4)),
]:
    pp = tf2.paragraphs[0] if txt.startswith("Cluster") else tf2.add_paragraph()
    pp.space_after = Pt(6)
    rr = pp.add_run(); rr.text = txt; _set(rr.font, sz, txt.startswith("Julia"), col)
notes(s, "Welcome. Our cluster looks at China strategies in the car industry. The central question is simple: how do you build a successful car company in China, and which ownership and financing model creates the most value over time? We compare three real companies, each standing for one model: SAIC, BYD, and Geely.")

# ============================ 1. INTRODUCTION ============================
content(2, "1. Introduction", "The Central Question",
    [(0,"How can a successful automotive company be built in China?"),
     (0,"Which governance and financing structure creates long-term shareholder value?"),
     (0,"Three models compared:"),
     (1,"Joint Venture → SAIC Motor"),
     (1,"Vertical Integration → BYD"),
     (1,"Cross-border M&A → Geely"),
     (0,"NIO’s capital-market model lies outside this team’s scope")],
    "Here is our central question. We do not just describe three companies, we compare them. Each company stands for one model. SAIC uses joint ventures with foreign partners. BYD builds everything itself. Geely grows by buying other companies. We left out NIO, the capital-market model, because it was not part of our team.")

content(3, "1. Introduction", "Why This Matters Now",
    [(0,"EVs reached about 55% of new car sales in China in 2025 (IEA, 2026)"),
     (0,"Chinese makers supplied around 60% of global EV sales that year"),
     (0,"The joint-venture model that defined China for three decades is unravelling"),
     (0,"EU countervailing duties turn governance into a trade-policy variable"),
     (0,"Central thesis: Governance structure, more than initial technology, determines long-term value creation")],
    "Why now? The market is changing very fast. In 2025 more than half of new cars sold in China were electric. Chinese brands now make about 60 percent of all electric vehicles in the world. The old joint-venture model is breaking down. Our main argument in one sentence: governance matters more than technology for long-term success.")

# ============================ 2. THEORY ============================
content(4, "2. Theory and Literature Review", "Governance and Ownership",
    [(0,"Shleifer & Vishny (1994): control rights, not cash-flow rights, drive resource allocation — politics enters the firm"),
     (0,"Megginson & Netter (2001): mixed-ownership firms underperform private peers"),
     (0,"Fan et al. (2007): politically connected Chinese firms underperform by ~18% post-IPO"),
     (0,"Jensen & Meckling (1976): concentrated ownership reduces the principal-agent conflict"),
     (0,"Allen et al. (2024): private Chinese firms structurally outperform state-controlled ones")],
    "Now the theory. This first group of papers is about ownership and control. Shleifer and Vishny show that whoever holds control rights steers the company, sometimes for political goals. Megginson and Netter, and Fan and co-authors, show that state-linked firms tend to perform worse. Jensen and Meckling explain why a strong single owner reduces conflicts. Allen and co-authors confirm that private firms do better in China.")

content(5, "2. Theory and Literature Review", "Capability and Strategy",
    [(0,"Holmes et al. (2015): “quid pro quo” — market access traded for technology transfer"),
     (0,"Bai et al. (2025): JV affiliation gave a +8.3% quality gain, but spillovers shrink as the gap closes"),
     (0,"Howell (2018): JV rents can suppress own innovation (“like opium”)"),
     (0,"Teece et al. (1997): dynamic capabilities — tacit, path-dependent, hard to imitate"),
     (0,"Luo & Tung (2007): “springboard” acquisitions; Inkpen & Beamish (1997): JVs turn unstable as learning becomes asymmetric")],
    "The second group of papers is about building capabilities. Holmes and co-authors describe the China deal: foreign firms get market access, China gets technology. Bai and co-authors measure a real quality gain from this. But Howell warns that easy joint-venture profits can kill the motivation to innovate. Teece explains why deep internal knowledge is hard to copy. Luo and Tung explain why firms like Geely buy foreign companies to catch up fast.")

# ============================ 3. DISCUSSION ============================
table_slide(6, "3. Discussion", "The Three Models at a Glance",
    [["", "SAIC Motor", "BYD", "Geely"],
     ["Model", "Joint Venture", "Vertical Integration", "M&A-driven"],
     ["Ownership", "State (SASAC ~73%)", "Founder (Wang ~17%)", "Founder (Li Shufu ~41%)"],
     ["Capital source", "JV dividends + state equity", "Operating cash flow", "Debt (acquisition-financed)"],
     ["Risk bearer", "State + JV partners (50:50)", "Equity holders (founder-led)", "Debt providers"],
     ["Capability mode", "Mandatory tech transfer", "Organic internal build", "Cross-border M&A"]],
    "This table sets up the whole comparison. Read it column by column. SAIC is state-owned and uses joint-venture money; the state and foreign partners share the risk. BYD is founder-led and pays for growth from its own cash, so the owners carry the risk. Geely is also founder-led but grows with debt to buy companies, so lenders carry much of the risk.",
    col_widths=[2.5, 3.2, 3.2, 3.2], fontsize=13)

content(7, "3. Discussion", "Governance: Who Controls, and What They Optimize",
    [(0,"SAIC: dual principal — SASAC (~73% state) plus a codified Party mandate → goals beyond shareholder value"),
     (0,"BYD: concentrated founder control (Wang ~17%, team >⅓) → fastest decision cycle, no principal-agent conflict"),
     (0,"Geely: founder apex (Li Shufu ~41%) → agile at the top, but multi-brand coordination cost"),
     (0,"Theory: Shleifer & Vishny (1994); Megginson & Netter (2001); Fan et al. (2007)")],
    "Governance first. SAIC has two bosses at once: the state holding company and the Communist Party. So decisions follow political goals too, not only profit. BYD has one strong founder, so it can decide the fastest. Geely also has a founder in control, but it owns many brands, so coordinating them costs time and effort.")

content(8, "3. Discussion", "Financing: Risk Allocation",
    [(0,"SAIC: capital-light JV — shared risk became a “capital trap” when SAIC-GM sales fell 56.5% in 2024"),
     (0,"BYD: self-financing — operating cash flow of RMB 133.5bn (2024); R&D spending exceeds net profit"),
     (0,"Geely: leverage-for-speed — gearing rose from 8.8% to 19.8% (FY2025); dependent on Volvo’s profit"),
     (0,"Same growth, three very different risk profiles")],
    "Financing next. SAIC's model looked cheap, because it split the cost with foreign partners. But when the GM joint venture collapsed in 2024, SAIC had to carry the losses without having control. BYD pays for everything from its own cash, so it is very independent. Geely borrows money to buy companies, which is fast but raises its debt.")

table_slide(9, "3. Discussion", "Comparative Financial Profile (FY2024 / FY2025)",
    [["Metric", "SAIC ’24", "SAIC ’25", "BYD ’24", "BYD ’25", "Geely ’24", "Geely ’25"],
     ["Revenue (RMB bn)", "627.6", "656.2", "777.1", "804.0", "240.2", "345.2"],
     ["Gross margin", "9.4%", "10.1%", "19.4%", "17.7%", "15.9%", "16.6%"],
     ["Net profit (RMB bn)", "1.7", "10.1", "40.3", "32.6", "16.8", "16.9"],
     ["Return on equity", "0.6%", "3.4%", "21.7%", "13.2%", "19.2%", "18.2%"],
     ["R&D (RMB bn)", "35.2", "33.6", "54.2", "63.4", "10.4", "17.6"],
     ["Op. cash flow (RMB bn)", "69.3", "34.3", "133.5", "59.1", "26.5", "47.3"],
     ["EU anti-subsidy tariff", "35.3%", "–", "17.0%", "–", "18.8%", "–"]],
    "This is the heart of the quantitative comparison. Look at return on equity: BYD made about 22 percent in 2024, Geely about 19, SAIC almost zero. Look at research spending: BYD spends more than its entire net profit. SAIC has the thinnest margins. The numbers tell the same story as the theory.",
    col_widths=[2.9, 1.55, 1.55, 1.55, 1.55, 1.55, 1.55], fontsize=11)

content(10, "3. Discussion", "Capability Building: Speed versus Depth",
    [(0,"SAIC: the JV “capability paradox” — a +8.3% quality gain (Bai et al., 2025), but no proprietary EV platforms"),
     (0,"Geely: acquisition turned into genuine co-development — the CMA platform, with measurable patent convergence"),
     (0,"BYD: organic build, slowest but deepest — non-imitable tacit knowledge (Teece et al., 1997)"),
     (0,"Trade-off: speed (Geely) vs depth (BYD) vs dependency (SAIC)")],
    "Now capabilities. SAIC learned to build good cars through its joint ventures, but it never built its own electric-car technology. Geely bought Volvo and turned it into real shared development, not one-way copying. BYD took the slow road and built everything itself, from batteries to software, which is the hardest for rivals to copy.")

content(11, "3. Discussion", "The EV Transition as a Structural Audit",
    [(0,"SAIC: most exposed — SAIC-GM at 37% capacity, VW MEB plant at 14%; a race between falling JV rents and a thin own-brand margin"),
     (0,"BYD: strongest — no combustion-engine legacy since 2022, 4 years global NEV leader, >1m exports in 2025 (+140%), 119 countries"),
     (0,"Geely: most asymmetric — 51.5% NEV share (H1 2025), +102% adjusted profit, but loss-making units (Polestar, Lotus)")],
    "The shift to electric cars works like an exam for each model. SAIC is under the most pressure: its joint-venture factories are almost empty. BYD is in the best shape: no old engine business, and four years as the world's electric-vehicle leader. Geely is in between: a strong core, but some loss-making brands create risk.")

stat_slide(12, "3. Discussion", "External Validation: The EU Tariff Differential",
    [("35.3%", "SAIC", ACC), ("17.0%", "BYD", BLUE), ("18.8%", "Geely", BLUE)],
    "An independent regulator priced SAIC’s state-dependence at more than double its private peers (European Commission, 2024).",
    "Here is outside proof for our argument. The European Union studied state subsidies and set import tariffs on Chinese electric cars. SAIC got 35 percent, more than double BYD and Geely. So an independent regulator confirmed the same point we found: SAIC depends most on the state. Governance even shows up in trade policy.")

# ============================ 4. LESSONS & CONCLUSIONS ============================
content(13, "4. Key Joint Lessons", "Lesson 1 — Mode of Capability Acquisition Determines Durability",
    [(0,"Three modes: mandatory transfer (SAIC), organic build (BYD), voluntary acquisition (Geely)"),
     (0,"Mandatory transfer buys short-run quality but displaces autonomous innovation"),
     (0,"Organic and voluntary modes preserve it — SAIC ended the cycle without proprietary EV platforms"),
     (0,"The trade-off across cycles: speed vs depth vs independence")],
    "Our first joint lesson. How you get your skills matters for the long run. Forced technology transfer gives quick quality, but it kills your own innovation. Building or buying voluntarily keeps it. There is always a trade-off between speed, depth, and independence.")

content(14, "4. Key Joint Lessons", "Lesson 2 — Governance Is a First-Order Determinant of Adaptability",
    [(0,"In fast markets, the speed of moving capital is itself a competitive advantage"),
     (0,"That speed depends on how control rights are distributed, not on resources"),
     (0,"Concentrated, commercial control adapts faster than political or JV-split control"),
     (0,"When change outpaces the decision cycle, governance becomes the binding constraint")],
    "Second lesson, and our main one. In fast-changing markets, being able to move money quickly is an advantage. That ability depends on who holds control. One commercial owner adapts faster than a mix of state, party, and joint-venture partners. When technology moves faster than the firm can decide, governance becomes the limiting factor.")

content(15, "4. Key Joint Lessons", "Lesson 3 — Capital Structure Determines Resilience",
    [(0,"Risk-sharing (SAIC) captures upside in growth but exposes the firm in downturns where it lacks control"),
     (0,"Self-financing (BYD) trades growth speed for resilience"),
     (0,"Leverage (Geely) trades balance-sheet safety for strategic speed"),
     (0,"Megginson & Netter (2001): mixed-ownership firms systematically underperform private peers")],
    "Third lesson. How you finance the firm decides how well you survive a downturn. Sharing risk is great while sales grow, but dangerous when they fall. Paying from your own cash is safer but slower. Borrowing is fast but risky. The benchmark study confirms that mixed state-private ownership underperforms.")

content(16, "4. Conclusions", "Answer to the Central Question",
    [(0,"BYD is the strongest model for long-term shareholder value"),
     (1,"Highest ROE in FY2024: 21.7% vs Geely 19.2% vs SAIC 0.6%"),
     (1,"Four consecutive years as global NEV market leader"),
     (1,"USD 5.6bn capital raise with sovereign wealth funds (2025)"),
     (0,"Not because vertical integration is inherently best — because it aligns commercial control with a resilient, self-financing structure"),
     (0,"Broader implication: when technology moves faster than firms can decide, governance is strategy")],
    "So what is the answer? BYD comes out strongest for long-term value. It has the best returns, four years as the world's electric-vehicle leader, and strong investor demand. The reason is not the factory model itself, but that one commercial owner plus self-financing makes it resilient when markets fall. The big takeaway: governance is not a small detail, governance is strategy.")

# ============================ 5. REFERENCES ============================
content(17, "5. References", "Selected Sources",
    [(0,"Theory: Shleifer & Vishny (1994); Megginson & Netter (2001); Fan et al. (2007); Jensen & Meckling (1976); Allen et al. (2024)"),
     (0,"Capability: Holmes et al. (2015); Bai et al. (2025); Howell (2018); Teece et al. (1997); Luo & Tung (2007); Inkpen & Beamish (1997); Zheng et al. (2022); Konda et al. (2022)"),
     (0,"Primary data: SAIC, BYD and Geely audited annual reports 2024–2025"),
     (0,"Regulatory & market: EU Implementing Regulation 2024/2754; IEA Global EV Outlook 2026"),
     (0,"Full reference list in the term paper (APA 7th edition)")],
    "These are our main sources. We used peer-reviewed academic papers for the theory, and primary data from the audited annual reports of all three companies, plus the EU regulation and the IEA outlook. The complete list, in APA style, is in our written paper.")

# ============================ 6. AI HIGHLIGHTS ============================
content(18, "6. AI Highlights", "Architecture and Workflow",
    [(0,"Researched with Claude Code (Anthropic): claude-sonnet-4-6 (primary) + claude-opus-4 (cross-case analysis)"),
     (0,"Not one chatbot, but a layer of specialized agents:"),
     (1,"wiki-ingest, wiki-query, wiki-lint, verify-claims, find-references"),
     (0,"Shared git repository, one wiki folder per case (pull before, push after)"),
     (0,"Source-quality rating (Q1–Q5 journal tiers × R1–R4 relevance) and a per-case writing guide")],
    "Now the AI part, which the seminar asks us to cover. We used Claude Code. We did not use one single chatbot; we used a system of special agents: one to read sources, one to answer questions, one to check the wiki, one to verify facts, one to find references. We worked in a shared git repository and rated every source by quality.")

content(19, "6. AI Highlights", "Was It Helpful? Value, Limits, and Cost",
    [(0,"Helpful: caught 3 real data errors before submission — EU tariff 36.3%→35.3%, a 10× currency error, a wrong-share-count EPS; all traced to audited reports"),
     (0,"Challenges: no paywalled-journal access; Geely PDFs not machine-readable; occasional currency/citation inconsistency; some wrong DOIs fixed by hand"),
     (0,"Setup: harder than a chatbot — learning the agent and wiki workflow took several hours"),
     (0,"Cost: Claude Pro ≈ €21 per person per month; no extra per-query charge")],
    "Was it helpful? Yes. It caught three real number mistakes before we handed in the paper, and it was great for handling many sources at once. But there were limits: it could not open paywalled journals, it could not read Geely's PDF report, and it sometimes mixed currencies or produced a wrong DOI. The setup took a few hours to learn, more than a normal chatbot. The cost was about 21 euros per person per month. Importantly, all the analysis and judgment stayed with us.")

# ============================ 20 — CLOSING ============================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.9), SW, Inches(0.06), ACC)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Thank you"; _set(r.font, 44, True, WHITE)
tb2 = s.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11.5), Inches(1.6))
tf2 = tb2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]
r2 = p2.add_run(); r2.text = "Governance, more than technology, determines who wins the EV transition."
_set(r2.font, 20, False, RGBColor(0xCB,0xD6,0xE6), italic=True)
p3 = tf2.add_paragraph(); p3.space_before = Pt(16)
r3 = p3.add_run(); r3.text = "Questions & Discussion"; _set(r3.font, 18, True, ACC)
notes(s, "Thank you for listening. To close with one line: technology matters, but governance decides who can actually use it fast enough. We are now happy to take your questions.")

out = r"C:\Users\User\CHINA-seminar-paper\research\presentation\china-strategies-julia.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
